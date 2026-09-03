#!/usr/bin/env python3
"""Design-guideline validator for NBG-design PPTX decks.

Checks the delivered .pptx against the nbg-design guardrails (palette, Aptos, real logo
lockups, native objects only, no collisions, footer clearance, legibility, closing
slide) and prints a FAIL / WARN / PASS report. Exit 1 on any FAIL.

Usage:
    source .venv/bin/activate
    python validate_deck.py <deck>.pptx [--json]
"""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

sys.path.insert(0, str(Path(__file__).resolve().parent))
from nbg_pptx import theme  # noqa: E402
from nbg_pptx.primitives import text_width_em  # noqa: E402

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}
THANKS = re.compile(r"thank\s*you|ευχαριστ|questions\?|q\s*&\s*a", re.I)


@dataclass
class Report:
    fails: list[str] = field(default_factory=list)
    warns: list[str] = field(default_factory=list)
    passes: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.fails

    def render(self) -> str:
        lines = ["NBG design-guideline validation", "=" * 40]
        for p in self.passes:
            lines.append(f"  PASS  {p}")
        for w in self.warns:
            lines.append(f"  WARN  {w}")
        for f in self.fails:
            lines.append(f"  FAIL  {f}")
        lines.append("-" * 40)
        lines.append(f"{len(self.passes)} pass, {len(self.warns)} warn, {len(self.fails)} fail -> {'PASS' if self.ok else 'FAIL'}")
        return "\n".join(lines)


@dataclass
class Box:
    name: str
    x: float
    y: float
    w: float
    h: float
    kind: str  # text | picture | shape | chart | table
    descr: str = ""
    runs: list = field(default_factory=list)  # (text, font, size_pt, color, alpha)
    paragraphs: list = field(default_factory=list)  # (runs[(text, size_px)], max_size_px, line_spacing)

    @property
    def x2(self):
        return self.x + self.w

    @property
    def y2(self):
        return self.y + self.h

    def intersects(self, o: "Box", tol: float = 1.0) -> bool:
        return not (self.x2 <= o.x + tol or o.x2 <= self.x + tol or self.y2 <= o.y + tol or o.y2 <= self.y + tol)

    @property
    def area(self):
        return self.w * self.h


def _px(emu) -> float:
    return (emu or 0) / theme.EMU_PER_PX


def _shape_boxes(slide) -> list[Box]:
    boxes = []
    for shp in slide.shapes:
        kind = "shape"
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            kind = "picture"
        elif shp.has_text_frame and shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            kind = "text"
        elif getattr(shp, "has_chart", False) and shp.has_chart:
            kind = "chart"
        elif getattr(shp, "has_table", False) and shp.has_table:
            kind = "table"
        b = Box(shp.name or "", _px(shp.left), _px(shp.top), _px(shp.width), _px(shp.height), kind)
        if kind == "picture":
            b.descr = shp._element.nvPicPr.cNvPr.get("descr", "")
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                ptext = "".join(r.text for r in p.runs)
                max_px = 0.0
                prun = []
                for r in p.runs:
                    r_pr = r._r.find(qn("a:rPr"))
                    size_pt = r.font.size.pt if r.font.size else None
                    color = alpha = None
                    if r_pr is not None:
                        sf = r_pr.find(qn("a:solidFill"))
                        if sf is not None and sf.find(qn("a:srgbClr")) is not None:
                            srgb = sf.find(qn("a:srgbClr"))
                            color = srgb.get("val")
                            a = srgb.find(qn("a:alpha"))
                            alpha = int(a.get("val")) / 100000 if a is not None else 1.0
                    b.runs.append((r.text, r.font.name, size_pt, color, alpha))
                    if size_pt:
                        max_px = max(max_px, size_pt / theme.PT_PER_PX)
                        prun.append((r.text, size_pt / theme.PT_PER_PX))
                lsp = p.line_spacing
                if lsp is None:
                    line_px = max_px * 1.2
                elif isinstance(lsp, float):
                    line_px = max_px * 1.2 * lsp
                else:  # exact Length
                    line_px = lsp.pt / theme.PT_PER_PX
                if ptext.strip() or p.runs:
                    b.paragraphs.append((prun, max_px, line_px))
        boxes.append(b)
    return boxes


def _slide_bg(slide) -> str | None:
    bg = slide._element.find(qn("p:cSld")).find(qn("p:bg"))
    if bg is None:
        return None
    srgb = bg.find(".//" + qn("a:srgbClr"))
    return srgb.get("val") if srgb is not None else None


def _blend(fg: str, alpha: float, bg: str) -> str:
    out = ""
    for i in (0, 2, 4):
        f, b = int(fg[i : i + 2], 16), int(bg[i : i + 2], 16)
        out += f"{int(round(alpha * f + (1 - alpha) * b)):02X}"
    return out


def validate(path: Path) -> Report:
    rep = Report()
    prs = Presentation(str(path))
    if (prs.slide_width, prs.slide_height) != (theme.SLIDE_W_EMU, theme.SLIDE_H_EMU):
        rep.fails.append(f"slide size is {prs.slide_width}x{prs.slide_height} EMU, expected 13.333in x 7.5in (16:9 = the 1920x1080 artboard)")
    else:
        rep.passes.append("16:9 slide size matches the 1920x1080 artboard")

    cat = prs.core_properties.category or ""
    style = cat.split(":", 1)[1] if cat.startswith("nbg-design-pptx:") else "design-system"
    lean = not style.endswith(";lean=off")
    style = style.split(";", 1)[0]
    briefing = style == "briefing"
    rep.passes.append(f"style: {style}" + ("" if lean else " (full-length deck: deck.lean is false, lean limits are advisory)"))

    z = zipfile.ZipFile(path)
    # --- palette + fonts from raw XML (slides and charts) ---
    bad_colors, bad_fonts = set(), set()
    base_palette = set(theme.BRIEF_PALETTE) if briefing else set(theme.PALETTE)
    for name in z.namelist():
        if not (name.startswith("ppt/slides/slide") or name.startswith("ppt/charts/chart")):
            continue
        root = etree.fromstring(z.read(name))
        allowed = set(base_palette)
        if name.startswith("ppt/charts/"):
            allowed |= theme.PALETTE_CHART_EXTRA | {"5A5F5A"}
        else:
            allowed |= {"5A5F5A"}
        for el in root.iter(qn("a:srgbClr")):
            v = (el.get("val") or "").upper()
            if v and v not in allowed:
                bad_colors.add(v)
        for tag in ("a:latin", "a:ea", "a:cs"):
            for el in root.iter(qn(tag)):
                tf = el.get("typeface") or ""
                if tf and not tf.startswith("+") and tf not in theme.ALLOWED_FONTS:
                    bad_fonts.add(tf)
    if bad_colors:
        rep.fails.append(f"colours outside the {'briefing' if briefing else 'NBG'} palette: {sorted(bad_colors)}")
    else:
        rep.passes.append("every colour is a palette token" + (" (briefing: white ground, #007B85 accent, grey ink, status tones)" if briefing else ""))
    if bad_fonts:
        rep.fails.append(f"fonts outside the Aptos family: {sorted(bad_fonts)}")
    else:
        rep.passes.append("Aptos family only (Light / Regular / SemiBold / Bold)")

    n = len(prs.slides)
    slide_area = theme.ARTBOARD_W * theme.ARTBOARD_H
    logo_ok = raster_ok = bounds_ok = overlap_ok = footer_ok = fonts_ok = overflow_ok = contrast_ok = True
    accent_warn = []
    white_ok = True
    for idx, slide in enumerate(prs.slides, start=1):
        boxes = _shape_boxes(slide)
        bg = _slide_bg(slide) or theme.PAPER
        dark = theme.is_dark(bg)
        tname = slide.name or ""
        if briefing and bg != "FFFFFF":
            rep.fails.append(f"slide {idx}: background #{bg}; briefing slides are always pure white")
            white_ok = False
        # logo
        logos = [b for b in boxes if b.descr.startswith("nbg-logo:")]
        if not logos:
            rep.fails.append(f"slide {idx}: no NBG logo lockup image (a real picture, never text or a drawn mark)")
            logo_ok = False
        else:
            for lg in logos:
                variant = lg.descr.split(":", 1)[1]
                if dark and variant != "knockout":
                    rep.fails.append(f"slide {idx}: '{variant}' logo on a dark surface; dark surfaces take the knockout lockup")
                    logo_ok = False
                if not dark and variant == "knockout":
                    rep.fails.append(f"slide {idx}: knockout (white) logo on a light surface")
                    logo_ok = False
                if lg.w / lg.h < 4.0 or lg.w / lg.h > 5.2:
                    rep.fails.append(f"slide {idx}: logo aspect ratio {lg.w / lg.h:.2f} is distorted")
                    logo_ok = False
        # rasters
        for b in boxes:
            if b.kind == "picture" and not (b.descr.startswith("nbg-photo:") or b.descr.startswith("user-photo:") or b.descr.startswith("nbg-logo:")):
                if b.area > 0.10 * slide_area:
                    rep.fails.append(f"slide {idx}: picture '{b.name}' covers {100 * b.area / slide_area:.0f}% of the slide and is not a declared photo asset (screenshot / rasterised content is forbidden)")
                    raster_ok = False
        # bounds
        for b in boxes:
            if b.x < -1 or b.y < -1 or b.x2 > theme.ARTBOARD_W + 1 or b.y2 > theme.ARTBOARD_H + 1:
                rep.fails.append(f"slide {idx}: '{b.name}' extends off the slide ({b.x:.0f},{b.y:.0f})-({b.x2:.0f},{b.y2:.0f})")
                bounds_ok = False
        # collisions: text vs text, text vs photo (footer exempt over photo; deco exempt)
        texts = [b for b in boxes if b.name.startswith("text:") and b.paragraphs]
        photos = [b for b in boxes if b.kind == "picture" and not b.descr.startswith("nbg-logo:")]
        for a in texts:
            for lg in logos:
                if a.intersects(lg):
                    rep.fails.append(f"slide {idx}: text '{a.name}' collides with the logo")
                    overlap_ok = False
        tables_charts = [b for b in boxes if b.kind in ("chart", "table")]
        for i, a in enumerate(texts):
            for bb in texts[i + 1 :]:
                if a.intersects(bb):
                    rep.fails.append(f"slide {idx}: text '{a.name}' overlaps text '{bb.name}'")
                    overlap_ok = False
            for ph in photos:
                if a.intersects(ph):
                    rep.fails.append(f"slide {idx}: text '{a.name}' overlaps photo '{ph.name}'")
                    overlap_ok = False
            for tc in tables_charts:
                if a.intersects(tc):
                    rep.fails.append(f"slide {idx}: text '{a.name}' overlaps {tc.kind} '{tc.name}'")
                    overlap_ok = False
        # footer clearance
        footer_top = theme.ARTBOARD_H - theme.FOOTER["bottom"] - theme.FOOTER["logo_h"]
        clearance = theme.BRIEF_FOOTER_CLEARANCE if briefing else theme.FOOTER_CLEARANCE
        has_footer = any(b.name.startswith("footer:") for b in boxes)
        if has_footer:
            for b in texts + tables_charts:
                if b.name == "text:footnote" and not briefing:
                    continue  # the design pins the footnote 46px above the footer band
                if b.y2 > footer_top - clearance + 1:
                    rep.fails.append(f"slide {idx}: '{b.name}' ends at y={b.y2:.0f}px, inside the {clearance}px footer clearance (footer at y={footer_top})")
                    footer_ok = False
        # font floor, overflow, contrast
        accents_used = set()
        for b in boxes:
            if not b.paragraphs:
                continue
            small_ok = b.name.startswith(("footer:", "text:eyebrow", "text:footnote", "text:meta", "pill:", "deco:")) or (
                briefing and b.name.startswith(("text:rail", "text:flow-label", "text:step", "text:matrix-h", "text:card", "text:band", "text:ask-title", "text:note", "text:box-title", "text:section", "text:phase"))
            )
            for text, font, size_pt, color, alpha in b.runs:
                if size_pt is None:
                    continue
                if size_pt < 7:
                    rep.fails.append(f"slide {idx}: '{b.name}' has {size_pt:.1f}pt text (floor is 7pt = 14px footnote)")
                    fonts_ok = False
                elif size_pt < 9 and not small_ok:
                    rep.warns.append(f"slide {idx}: '{b.name}' has {size_pt:.1f}pt text; body copy should stay at or above 11pt (22px)")
                if color in (theme.TEAL, theme.BRIGHT, theme.CYAN):
                    accents_used.add(color)
                if color and b.name.startswith("text:") and not any(b.intersects(d) for d in boxes if d.name.startswith(("deco:", "bg:", "photo:"))):
                    eff = _blend(color, alpha if alpha is not None else 1.0, bg)
                    cr = theme.contrast_ratio(eff, bg)
                    muted = alpha is not None and alpha <= 0.6  # deliberately muted secondary copy (bilingual Greek line)
                    if cr < (1.5 if muted else 2.0):
                        rep.fails.append(f"slide {idx}: '{b.name}' text #{color} (alpha {alpha:.2f}) has contrast {cr:.1f}:1 on #{bg}")
                        contrast_ok = False
                    elif cr < 2.4 and size_pt < 14 and (alpha is None or alpha > 0.6):
                        rep.warns.append(f"slide {idx}: '{b.name}' small text has low contrast {cr:.1f}:1 on #{bg}")
            if b.name.startswith("text:") or b.name.startswith("footer:"):
                need = 0.0
                for prun, size_px, line_px in b.paragraphs:
                    if not size_px:
                        continue
                    # advance width per run, honouring explicit line breaks
                    lines, width = 1, 0.0
                    for text, rs in prun:
                        for k, seg in enumerate(text.split("\n")):
                            if k:
                                lines += max(1, math.ceil(width / b.w)) if width else 1
                                width = 0.0
                            width += text_width_em(seg) * rs
                    lines_total = (lines - 1) + max(1, math.ceil(width / b.w)) if width else lines
                    need += lines_total * line_px
                if need > b.h * 1.5 + 4:
                    rep.fails.append(f"slide {idx}: '{b.name}' cannot fit its zone (~{need:.0f}px of text in {b.h:.0f}px); shorten the copy, split the slide, or drop to a single language")
                    overflow_ok = False
                elif need > b.h * 1.08 + 4:
                    rep.warns.append(f"slide {idx}: '{b.name}' may overflow its box (~{need:.0f}px of text in {b.h:.0f}px); check the render")
                    overflow_ok = False
        if len(accents_used) > 1:
            accent_warn.append(f"slide {idx}: {len(accents_used)} accent colours in text ({sorted(accents_used)}); the system keeps one accent per slide")
        # closing slide
        if idx == n:
            alltext = " ".join(t for b in boxes for t, *_ in b.runs)
            if THANKS.search(alltext):
                msg = "last slide carries a 'Thank you' / 'Questions' message; " + ("a briefing ends on its asks / next steps" if briefing else "close with the plain back cover instead")
                (rep.warns if (briefing and not lean) else rep.fails).append(msg)
            if not briefing and tname != "back_cover":
                rep.warns.append("last slide is not the `back_cover` template")
        # em-dash advisory
        for b in boxes:
            for t, *_ in b.runs:
                if "—" in t:
                    rep.warns.append(f"slide {idx}: em-dash in '{b.name}' (the studio guideline avoids em-dashes in slide copy; the design system uses a middle dot '·')")
                    break

    for msg, ok in (
        ("real NBG lockup on every slide, correct variant for the surface", logo_ok),
        ("no rasterised content - only declared photo assets and logos", raster_ok),
        ("every object inside the 1920x1080 artboard", bounds_ok),
        ("no collisions between text blocks, photos, charts or tables", overlap_ok),
        (f"{theme.BRIEF_FOOTER_CLEARANCE if briefing else theme.FOOTER_CLEARANCE}px clearance above the footer band", footer_ok),
        ("no text below the 7pt floor", fonts_ok),
        ("effective text contrast at or above 2:1", contrast_ok),
    ):
        if ok:
            rep.passes.append(msg)
    if overflow_ok:
        rep.passes.append("no estimated text overflow")
    if briefing:
        if white_ok:
            rep.passes.append("pure white ground on every slide")
        if n > theme.BRIEF_LEAN_FAIL and lean:
            rep.fails.append(f"{n} slides: a briefing deck must stay lean (hard limit {theme.BRIEF_LEAN_FAIL}); merge or drop slides, move detail to notes")
        elif n > theme.BRIEF_LEAN_FAIL:
            rep.warns.append(f"{n} slides: full-length deck by request (deck.lean: false); the lean target for briefings is {theme.BRIEF_LEAN_WARN}")
        elif n > theme.BRIEF_LEAN_WARN:
            rep.warns.append(f"{n} slides: lean target is {theme.BRIEF_LEAN_WARN} or fewer for a briefing")
        else:
            rep.passes.append(f"lean: {n} slides (target <= {theme.BRIEF_LEAN_WARN})")
    rep.warns.extend(accent_warn if not briefing else [])
    rep.warns = list(dict.fromkeys(rep.warns))
    rep.fails = list(dict.fromkeys(rep.fails))
    rep.passes.append(f"{n} slides")
    return rep


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck", type=Path)
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()
    rep = validate(args.deck)
    if args.json:
        print(json.dumps({"ok": rep.ok, "fails": rep.fails, "warns": rep.warns, "passes": rep.passes}, indent=2, ensure_ascii=False))
    else:
        print(rep.render())
    return 0 if rep.ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
