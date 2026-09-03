"""Low-level python-pptx drawing primitives that reproduce the CSS used by the
NBG HTML slide templates as *native, editable* PowerPoint objects.

Every helper takes artboard pixels (1920x1080) and writes real shapes, text boxes,
pictures, gradient/pattern fills, alpha and effects - never rasterised content.
Shapes are named with a role prefix (`bg:` `deco:` `photo:` `text:` `logo:` `footer:`)
so the validator can reason about layering and collisions.
"""

from __future__ import annotations

import math
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from . import theme
from .theme import px

_FILL_TAGS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill")

# ---------------------------------------------------------------------------
# XML helpers
# ---------------------------------------------------------------------------


def _remove_fills(sp_pr) -> None:
    for tag in _FILL_TAGS:
        for el in sp_pr.findall(qn(tag)):
            sp_pr.remove(el)


def _insert_fill(sp_pr, fill_el) -> None:
    """Insert a fill element in the schema-correct position (after the geometry)."""
    _remove_fills(sp_pr)
    geom = sp_pr.find(qn("a:prstGeom"))
    if geom is None:
        geom = sp_pr.find(qn("a:custGeom"))
    if geom is not None:
        geom.addnext(fill_el)
    else:
        xfrm = sp_pr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.addnext(fill_el)
        else:
            sp_pr.insert(0, fill_el)


def _srgb(hex6: str, alpha: float | None = None):
    el = etree.SubElement(etree.Element("dummy"), qn("a:srgbClr"))
    el.set("val", theme.normalize_hex(hex6))
    if alpha is not None and alpha < 1.0:
        a = etree.SubElement(el, qn("a:alpha"))
        a.set("val", str(int(round(max(0.0, min(1.0, alpha)) * 100000))))
    return el


def _solid_fill_el(hex6: str, alpha: float | None = None):
    sf = etree.Element(qn("a:solidFill"))
    sf.append(_srgb(hex6, alpha))
    return sf


def _no_line(shape) -> None:
    shape.line.fill.background()


def _no_shadow(shape) -> None:
    # Writes an empty <a:effectLst/> so the theme's default shadow is not inherited.
    shape.shadow.inherit = False


def _css_angle_to_ooxml(css_deg: float) -> int:
    """CSS linear-gradient angle (0 = to top, clockwise) -> OOXML `ang`
    (0 = to right, clockwise, in 60000ths of a degree)."""
    return int(round(((css_deg - 90) % 360) * 60000))


# ---------------------------------------------------------------------------
# Shapes
# ---------------------------------------------------------------------------


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None,
    alpha: float | None = None,
    radius: float = 0,
    round_corner: str | None = None,
    name: str = "deco:rect",
    line: str | None = None,
    line_alpha: float | None = None,
    line_px: float = 1.0,
):
    """Rectangle with optional uniform corner radius (px) or a single rounded
    corner (`round_corner` = 'tl' | 'tr' | 'bl' | 'br') and an optional outline."""
    if round_corner:
        shape_type = MSO_SHAPE.ROUND_1_RECTANGLE
    elif radius > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    shp.name = name
    if radius > 0:
        shp.adjustments[0] = min(0.5, radius / max(1.0, min(w, h)))
    if round_corner:
        # round1Rect rounds the top-right corner; flip to move it.
        xfrm = shp._element.spPr.find(qn("a:xfrm"))
        if round_corner in ("tl", "bl"):
            xfrm.set("flipH", "1")
        if round_corner in ("bl", "br"):
            xfrm.set("flipV", "1")
    if fill is None:
        shp.fill.background()
    else:
        _insert_fill(shp._element.spPr, _solid_fill_el(fill, alpha))
    if line is None:
        _no_line(shp)
    else:
        _set_outline(shp, line, line_alpha, line_px)
    _no_shadow(shp)
    # Autoshapes carry an empty text body; keep it empty and non-wrapping.
    shp.text_frame.text = ""
    return shp


def _set_outline(shp, color: str, alpha: float | None, width_px: float) -> None:
    ln = shp.line
    ln.width = Pt(width_px * theme.PT_PER_PX)
    ln.color.rgb = RGBColor.from_string(theme.normalize_hex(color))
    if alpha is not None and alpha < 1.0:
        srgb = shp._element.spPr.find(qn("a:ln")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))


def add_oval(slide, x: float, y: float, w: float, h: float, *, fill: str | None, alpha: float | None = None, line: str | None = None, line_px: float = 1.0, name: str = "deco:oval"):
    """Ellipse with optional fill and outline (the briefing cover mark, step numbers)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(w), px(h))
    shp.name = name
    if fill is None:
        shp.fill.background()
    else:
        _insert_fill(shp._element.spPr, _solid_fill_el(fill, alpha))
    if line is None:
        _no_line(shp)
    else:
        _set_outline(shp, line, None, line_px)
    _no_shadow(shp)
    shp.text_frame.text = ""
    return shp


def set_shape_text(shp, text: str, *, size: float, weight: int, color: str, align: str = "center", caps: bool = False, letter_spacing: float = 0, inset_x: float = 0, inset_y: float = 0):
    """Put a single run of text inside an autoshape (pills, badges, step numbers)."""
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = px(inset_x)
    tf.margin_top = tf.margin_bottom = px(inset_y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = _ALIGN[align]
    p.line_spacing = Pt(theme.pt(size) * 1.0)
    run = p.add_run()
    run.text = text
    _apply_run_style(run, size_px=size, weight=weight, color=color, alpha=None, italic=False, letter_spacing=letter_spacing, caps=caps)
    return shp


def add_gradient_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    stops: list[tuple[float, str, float | None]],
    css_angle: float,
    radius: float = 0,
    name: str = "deco:gradient",
):
    """Linear gradient rectangle. `stops` = [(position 0..1, hex, alpha|None)]."""
    shp = add_rect(slide, x, y, w, h, fill=None, radius=radius, name=name)
    grad = etree.Element(qn("a:gradFill"))
    grad.set("rotWithShape", "1")
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, hex6, alpha in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(int(round(pos * 100000))))
        gs.append(_srgb(hex6, alpha))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(_css_angle_to_ooxml(css_angle)))
    lin.set("scaled", "0")
    _insert_fill(shp._element.spPr, grad)
    return shp


def add_pattern_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    preset: str,
    fg: str,
    fg_alpha: float,
    bg: str,
    name: str = "deco:pattern",
):
    """Pattern-filled rectangle (e.g. the diagonal hairlines of Cover 3)."""
    shp = add_rect(slide, x, y, w, h, fill=None, name=name)
    patt = etree.Element(qn("a:pattFill"))
    patt.set("prst", preset)
    fg_el = etree.SubElement(patt, qn("a:fgClr"))
    fg_el.append(_srgb(fg, fg_alpha))
    bg_el = etree.SubElement(patt, qn("a:bgClr"))
    bg_el.append(_srgb(bg))
    _insert_fill(shp._element.spPr, patt)
    return shp


def add_hairline(slide, x: float, y: float, w: float, *, color: str, alpha: float | None, thickness: float = 1, name="deco:hairline"):
    return add_rect(slide, x, y, w, thickness, fill=color, alpha=alpha, name=name)


def add_diagonal_hairlines(slide, *, spacing: float, color: str, alpha: float, width_px: float = 0.5, name: str = "bg:pattern"):
    """The Cover 3 background: 45-degree hairlines every `spacing` px across the whole
    artboard, drawn as a single group of native line connectors (CSS `pattern` equivalent)."""
    from pptx.enum.shapes import MSO_CONNECTOR

    W, H = theme.ARTBOARD_W, theme.ARTBOARD_H
    group = slide.shapes.add_group_shape()
    group.name = name
    # lines run from bottom-left to top-right; sweep the start point along the left edge
    # and then along the bottom edge so the whole rectangle is covered.
    starts = [(0, y) for y in range(int(spacing), H + int(spacing), int(spacing))]
    starts += [(x, H) for x in range(int(spacing), W, int(spacing))]
    for sx, sy in starts:
        # from (sx, sy) travel up-right until leaving the artboard
        d = min(sy, W - sx)
        ex, ey = sx + d, sy - d
        ln = group.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(sx), px(sy), px(ex), px(ey))
        ln.name = "deco:hairline"
        ln.line.width = Pt(width_px * theme.PT_PER_PX)
        ln.line.color.rgb = RGBColor.from_string(color)
        srgb = ln._element.spPr.find(qn("a:ln")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))
    return group


def add_shadow(shape, *, blur_px: float, dist_px: float, css_deg: float, color: str, alpha: float) -> None:
    """Soft outer shadow (CSS box-shadow equivalent). css_deg 90 = downwards."""
    sp_pr = shape._element.spPr
    for el in sp_pr.findall(qn("a:effectLst")):
        sp_pr.remove(el)
    eff = etree.SubElement(sp_pr, qn("a:effectLst"))
    sh = etree.SubElement(eff, qn("a:outerShdw"))
    sh.set("blurRad", str(px(blur_px)))
    sh.set("dist", str(px(dist_px)))
    sh.set("dir", str(int(round((css_deg % 360) * 60000))))
    sh.set("algn", "ctr")
    sh.set("rotWithShape", "0")
    sh.append(_srgb(color, alpha))


# ---------------------------------------------------------------------------
# Pictures
# ---------------------------------------------------------------------------


def _image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        return im.size


def add_picture_cover(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    radius: float = 0,
    alpha: float | None = None,
    name: str = "photo:image",
    descr: str = "",
):
    """Place an image with CSS `object-fit: cover` semantics: the frame is exactly
    (w, h) and the source is centre-cropped, never distorted. Optional rounded
    corners and opacity are applied natively to the picture."""
    iw, ih = _image_size(path)
    scale = max(w / iw, h / ih)
    vis_w, vis_h = (w / scale) / iw, (h / scale) / ih
    pic = slide.shapes.add_picture(str(path), px(x), px(y), px(w), px(h))
    pic.name = name
    pic.crop_left = pic.crop_right = max(0.0, (1 - vis_w) / 2)
    pic.crop_top = pic.crop_bottom = max(0.0, (1 - vis_h) / 2)
    if radius > 0:
        pic.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        geom = pic._element.spPr.find(qn("a:prstGeom"))
        av = geom.find(qn("a:avLst"))
        if av is None:
            av = etree.SubElement(geom, qn("a:avLst"))
        for gd in list(av):
            av.remove(gd)
        gd = etree.SubElement(av, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {int(round(min(0.5, radius / max(1.0, min(w, h))) * 100000))}")
    if alpha is not None and alpha < 1.0:
        blip = pic._element.xpath(".//a:blip")[0]
        amf = etree.SubElement(blip, qn("a:alphaModFix"))
        amf.set("amt", str(int(round(alpha * 100000))))
    pic._element.nvPicPr.cNvPr.set("descr", descr or name)
    _no_line(pic)
    return pic


def add_logo(slide, variant: str, x: float, y: float, height: float, *, anchor: str = "tl", name: str | None = None):
    """Insert the real NBG lockup image, preserving its native aspect ratio.
    `anchor` 'tl' places the top-left at (x, y); 'bl' the bottom-left; 'c' the centre."""
    path = theme.LOGOS[variant]
    iw, ih = _image_size(path)
    w = height * iw / ih
    if anchor == "bl":
        y = y - height
    elif anchor == "c":
        x, y = x - w / 2, y - height / 2
    pic = slide.shapes.add_picture(str(path), px(x), px(y), px(w), px(height))
    pic.name = name or f"logo:{variant}"
    pic._element.nvPicPr.cNvPr.set("descr", f"nbg-logo:{variant}")
    _no_line(pic)
    return pic


# ---------------------------------------------------------------------------
# Text
# ---------------------------------------------------------------------------

_ALIGN = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _apply_run_style(run, *, size_px: float, weight: int, color: str, alpha: float | None, italic: bool, letter_spacing: float, caps: bool, baseline: int | None = None, lang: str | None = None):
    face, bold = theme.font_for_weight(weight)
    f = run.font
    f.name = face
    f.size = Pt(theme.pt(size_px))
    f.bold = bold
    f.italic = italic
    f.color.rgb = RGBColor.from_string(theme.normalize_hex(color))
    r_pr = run._r.get_or_add_rPr()
    if alpha is not None and alpha < 1.0:
        srgb = r_pr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))
    if letter_spacing:
        r_pr.set("spc", str(int(round(letter_spacing * theme.PT_PER_PX * 100))))
    if caps:
        r_pr.set("cap", "all")
    if baseline:
        r_pr.set("baseline", str(int(baseline)))
    if lang:
        r_pr.set("lang", lang)
    # Explicit Greek/Latin face so PowerPoint never substitutes the theme font.
    for tag in ("a:latin", "a:ea", "a:cs"):
        for el in r_pr.findall(qn(tag)):
            r_pr.remove(el)
    latin = etree.SubElement(r_pr, qn("a:latin"))
    latin.set("typeface", face)
    ea = etree.SubElement(r_pr, qn("a:ea"))
    ea.set("typeface", face)
    cs = etree.SubElement(r_pr, qn("a:cs"))
    cs.set("typeface", face)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    paragraphs: list[list[dict]],
    *,
    size: float,
    weight: int = 400,
    color: str,
    alpha: float | None = None,
    line_height: float = 1.2,
    letter_spacing: float = 0,
    align: str = "left",
    anchor: str = "top",
    caps: bool = False,
    para_space_after: float = 0,
    para_space_before: float = 0,
    name: str = "text:block",
    bullet: dict | None = None,
    numbered: bool = False,
):
    """Add a text box. `paragraphs` is a list of paragraphs; each paragraph is a
    list of span dicts: {text, italic?, weight?, color?, alpha?, size?, br?}.
    A span with `br: True` inserts a line break (semantic <br>) inside the paragraph.
    Text boxes have zero insets and no autofit, matching the HTML absolute layout."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _ANCHOR[anchor]
    first = True
    for spans in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = _ALIGN[align]
        # CSS line-height is a multiple of the font size; OOXML percentages multiply the
        # font's own line gap (~1.2 em for Aptos), so write exact points instead.
        p.line_spacing = Pt(theme.pt(size) * line_height)
        if para_space_after:
            p.space_after = Pt(theme.pt(para_space_after))
        if para_space_before:
            p.space_before = Pt(theme.pt(para_space_before))
        if bullet or numbered:
            _set_bullet(p, size, bullet or {}, numbered)
        for span in spans:
            if span.get("br"):
                p.add_line_break()
                continue
            run = p.add_run()
            run.text = span["text"]
            span_alpha = span.get("alpha", alpha)
            _apply_run_style(
                run,
                size_px=span.get("size", size),
                weight=span.get("weight", weight),
                color=span.get("color", color),
                alpha=span_alpha,
                italic=span.get("italic", False),
                letter_spacing=span.get("letter_spacing", letter_spacing),
                caps=span.get("caps", caps),
                baseline=span.get("baseline"),
                lang=span.get("lang"),
            )
    return tb


def _set_bullet(p, size_px: float, opts: dict, numbered: bool) -> None:
    """Native bullet / numbering on a paragraph: hanging indent, accent-coloured glyph."""
    p_pr = p._p.get_or_add_pPr()
    indent = float(opts.get("indent", size_px * 1.1))
    p_pr.set("marL", str(px(indent)))
    p_pr.set("indent", str(-px(indent)))
    for tag in ("a:buNone", "a:buClr", "a:buSzPct", "a:buFont", "a:buChar", "a:buAutoNum"):
        for el in p_pr.findall(qn(tag)):
            p_pr.remove(el)
    color = opts.get("color")
    if color:
        bu_clr = etree.SubElement(p_pr, qn("a:buClr"))
        bu_clr.append(_srgb(color))
    bu_sz = etree.SubElement(p_pr, qn("a:buSzPct"))
    bu_sz.set("val", str(int(opts.get("size_pct", 100 if numbered else 70) * 1000)))
    if numbered:
        num = etree.SubElement(p_pr, qn("a:buAutoNum"))
        num.set("type", "arabicPeriod")
    else:
        bu_font = etree.SubElement(p_pr, qn("a:buFont"))
        bu_font.set("typeface", "Aptos")
        bu_char = etree.SubElement(p_pr, qn("a:buChar"))
        bu_char.set("char", opts.get("char", "\u25cf"))


# ---------------------------------------------------------------------------
# Estimation helpers (HTML flows text; PPTX needs explicit boxes)
# ---------------------------------------------------------------------------

AVG_CHAR_EM = 0.52  # Aptos average advance width in em (Latin lowercase, mixed text)


def char_em(ch: str) -> float:
    """Approximate Aptos advance width of one character, in em."""
    o = ord(ch)
    if ch == " ":
        return 0.26
    if ch in ".,:;'|!":
        return 0.28
    if ch.isdigit():
        return 0.55
    if 0x0370 <= o <= 0x03FF:  # Greek: wider lowercase, and capitals are common
        return 0.66 if ch.isupper() else 0.58
    if ch.isupper():
        return 0.66
    if ch in "mwMW":
        return 0.85
    if ch in "iljtfr":
        return 0.30
    return AVG_CHAR_EM


def text_width_em(text: str) -> float:
    return sum(char_em(c) for c in text)


def estimate_lines(text: str, size_px: float, width_px: float, *, avg_em: float | None = None) -> int:
    """Line count for wrapped text (explicit \\n honoured), from per-character advances."""
    total = 0
    for raw in text.split("\n"):
        if not raw:
            total += 1
            continue
        w = (len(raw) * avg_em if avg_em else text_width_em(raw)) * size_px
        total += max(1, math.ceil(w / max(1.0, width_px)))
    return max(1, total)


def estimate_height(text: str, size_px: float, width_px: float, line_height: float) -> float:
    return estimate_lines(text, size_px, width_px) * size_px * line_height
