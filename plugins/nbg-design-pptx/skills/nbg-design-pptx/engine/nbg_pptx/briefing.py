"""Briefing-style templates: the internal team-briefing look of the reference deck
(`index 1.html`, "SOFIA Chat & Voice - Release Outlook"), rebuilt as native PowerPoint.

Visual rules (fixed, not configurable):
- pure white ground on every slide;
- ONE accent (#007B85) used only for highlights, rules, tags, bullets and card edges;
- grey ink hierarchy: ink #0E1B1D, muted #5B6B6D, faint #8A9A9C, hairlines #E2E8E9;
- semantic status tones (ok / warn / stop / wait) only inside pills, dots and callouts;
- Aptos: SemiBold for titles and strong spans, Regular for body;
- a fixed vertical rhythm (head at y=64, hairline at y=130, body from y=156 to y=944,
  footer band at y=1016) - templates never "adjust whitespace", they refuse copy that
  does not fit (SpecError), because a briefing deck has to stay lean.

Geometry is the reference's 1280x720 stage scaled by 1.5 onto the 1920x1080 artboard.
"""

from __future__ import annotations

from pptx.util import Pt

from . import theme
from .charts import add_chart
from .primitives import (
    _apply_run_style,
    add_logo,
    add_oval,
    add_rect,
    add_text,
    estimate_lines,
    set_shape_text,
    text_width_em,
)
from .templates import Ctx
from .text import SpecError, joined_plain, plain, resolve_copy, spans_for

B = theme.BRIEF
TONES = theme.BRIEF_TONES

W, H = theme.ARTBOARD_W, theme.ARTBOARD_H
MX = 88  # side margin (58px on the 1280 stage)
CW = W - 2 * MX  # 1744
TOP = 64
HEAD_RULE_Y = 130
BODY_TOP = 156
FOOTER_TOP = H - theme.FOOTER["bottom"] - theme.FOOTER["logo_h"]  # 1016
BODY_BOTTOM = FOOTER_TOP - theme.BRIEF_FOOTER_CLEARANCE  # 976
GAP = 24

BODY = 19  # body copy px (12.5px on the 1280 stage)
SMALL = 15  # labels / tags / rail notes px (10-10.5px on the stage)


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _err(ctx: Ctx, msg: str) -> SpecError:
    return SpecError(f"slide {ctx.index + 1} ({ctx.template}): {msg}")


def _tone(ctx: Ctx, name: str | None, *, default: str = "accent") -> dict:
    key = default if name is None else str(name)
    if key not in TONES:
        raise _err(ctx, f"tone {key!r} is not one of {sorted(TONES)}")
    return TONES[key]


def _variants(ctx: Ctx, value, path: str) -> list[tuple[str, bool]]:
    return resolve_copy(value, ctx.lang, path=f"slide {ctx.index + 1}.{path}")


def _width(variants, size: float, weight: int = 400) -> float:
    k = 1.15 if weight >= 600 else 1.0
    return max((text_width_em(plain(t)) for t, _ in variants), default=0.0) * size * k


def _est(variants, size: float, width: float, lh: float, weight: int = 400) -> float:
    k = 1.15 if weight >= 600 else 1.0
    total = 0.0
    for text, _ in variants:
        total += estimate_lines(plain(text), size * k, width * 0.9) * size * lh
    return total


def _paras(ctx: Ctx, variants, *, color: str, weight: int = 400, alpha=None) -> list[list[dict]]:
    return spans_for(
        variants,
        base_color=color,
        accent_color=B["accent"],
        base_weight=weight,
        base_alpha=alpha,
        emphasis_weight=weight,
        strong_color=B["ink"],
        lang_hint=ctx.lang,
    )


def _tx(ctx: Ctx, x, y, w, h, variants, *, size, weight=400, color, alpha=None, lh=1.4, ls=0.0, align="left", anchor="top", caps=False, name, para_after=0, bullet=None, numbered=False):
    if not variants:
        return None
    return add_text(ctx.slide, x, y, w, h, _paras(ctx, variants, color=color, weight=weight, alpha=alpha), size=size, weight=weight, color=color, alpha=alpha, line_height=lh, letter_spacing=ls, align=align, anchor=anchor, caps=caps, para_space_after=para_after, name=name, bullet=bullet, numbered=numbered)


def _items(ctx: Ctx, items, path: str) -> list[list[tuple[str, bool]]]:
    if items is None:
        return []
    if not isinstance(items, list):
        raise _err(ctx, f"`{path}` must be a list")
    return [_variants(ctx, it, f"{path}[{i}]") for i, it in enumerate(items)]


def _list_block(ctx: Ctx, x, y, w, items: list, *, size=BODY, lh=1.6, color=None, numbered=False, name="text:list", gap=6) -> float:
    """Bulleted (accent dot) or numbered list, one paragraph per item. Returns the height."""
    if not items:
        return 0.0
    paragraphs: list[list[dict]] = []
    est = 0.0
    for it in items:
        paragraphs.extend(_paras(ctx, it, color=color or B["muted"]))
        est += _est(it, size, w - size * 1.1, lh) + gap
    est -= gap
    bullet = {"indent": size * 1.9} if numbered else {"color": B["accent"], "char": "●", "size_pct": 45, "indent": size * 1.1}
    add_text(ctx.slide, x, y, w, est, paragraphs, size=size, color=color or B["muted"], line_height=lh, para_space_after=gap, name=name, bullet=bullet if not numbered else None, numbered=numbered)
    return est


def _shape_runs(shp, runs: list[tuple[str, float, int, str]], *, align="left", inset_l=0, inset_r=0) -> None:
    """Several styled runs inside one autoshape: [(text, size_px, weight, color)]."""
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left, tf.margin_right = theme.px(inset_l), theme.px(inset_r)
    tf.margin_top = tf.margin_bottom = 0
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    for text, size, weight, color in runs:
        r = p.add_run()
        r.text = text
        _apply_run_style(r, size_px=size, weight=weight, color=color, alpha=None, italic=False, letter_spacing=0, caps=False)
    p.line_spacing = Pt(theme.pt(max(s for _, s, _, _ in runs)))


def _pill(ctx: Ctx, x, y, label: str, tone: dict, *, h=30, size=15.5, name="pill:status") -> float:
    """Status pill: tinted rounded capsule, small dot, SemiBold label. Returns the width."""
    w = text_width_em(label) * size * 1.2 + 14 + 13 + 8 + 12
    shp = add_rect(ctx.slide, x, y, w, h, fill=tone["bg"], radius=h / 2, name=name)
    add_oval(ctx.slide, x + 13, y + h / 2 - 4, 8, 8, fill=tone["fg"], name="deco:pill-dot")
    _shape_runs(shp, [(label, size, 600, tone["fg"])], inset_l=27, inset_r=13)
    return w


def _tag(ctx: Ctx, x, y, label: str, *, h=30, size=SMALL) -> float:
    w = text_width_em(label.upper()) * size * 1.1 + 22
    shp = add_rect(ctx.slide, x, y, w, h, fill=B["accent_xlt"], radius=6, line=B["accent"], line_px=1.5, name="deco:tag")
    set_shape_text(shp, label, size=size, weight=700, color=B["accent"], caps=True, letter_spacing=1.0)
    return w


def _badge(ctx: Ctx, x, y, label: str, tone: dict, *, h=34, size=16) -> float:
    w = text_width_em(label.upper()) * size * 1.15 + 30
    shp = add_rect(ctx.slide, x, y, w, h, fill=tone["fg"], radius=7, name="deco:badge")
    set_shape_text(shp, label, size=size, weight=700, color=B["white"], caps=True, letter_spacing=1.6)
    return w


def _kicker(ctx: Ctx, x, y, w, variants, *, name="text:eyebrow", size=16, color=None) -> float:
    """Small uppercase accent label (card headings, panel titles, cover eyebrow)."""
    _tx(ctx, x, y, w, 24, variants, size=size, weight=600, color=color or B["accent"], caps=True, ls=size * 0.12, lh=1.2, name=name)
    return 24


def _step_circle(ctx: Ctx, x, y, n: int, *, cond=False, d=28) -> None:
    if cond:
        shp = add_oval(ctx.slide, x, y, d, d, fill=B["white"], line=B["accent"], line_px=1.5, name="pill:step")
        set_shape_text(shp, str(n), size=14, weight=700, color=B["accent"])
    else:
        shp = add_oval(ctx.slide, x, y, d, d, fill=B["accent"], name="pill:step")
        set_shape_text(shp, str(n), size=14, weight=700, color=B["white"])


# ---------------------------------------------------------------------------
# shared chrome: head, footer, progress bar
# ---------------------------------------------------------------------------


def _head(ctx: Ctx, s: dict) -> float:
    """Slide head: [tag] title [sub] ......... meta (right), hairline. Returns the body top."""
    x, y = MX, TOP
    tag = ctx.copy(s, "tag")
    if tag:
        x += _tag(ctx, x, y + 9, joined_plain(tag)) + 16
    title = ctx.copy(s, "title", required=True)
    sub = ctx.copy(s, "sub")
    meta = s.get("meta")
    meta_w = 540 if meta else 0
    tw = MX + CW - x - (meta_w + 24 if meta else 0)
    paragraphs = _paras(ctx, title, color=B["ink"], weight=600)
    est_w = _width(title, 34, 600)
    if sub:
        sp = _paras(ctx, sub, color=B["faint"], weight=400)
        paragraphs[-1].append({"text": "  \u2013 ", "color": B["faint"], "weight": 400})
        paragraphs[-1].extend(sp[0])
        est_w += _width(sub, 34) + 34 * 1.2
    if est_w * 1.15 > tw:
        raise _err(ctx, f"the head is one line: title + sub measure ~{est_w * 1.15:.0f}px for a {tw:.0f}px line; shorten the title or the sub, or drop `meta`")
    add_text(ctx.slide, x, y, tw, 50, paragraphs, size=34, weight=600, color=B["ink"], line_height=1.2, letter_spacing=-0.4, anchor="bottom", name="text:title")
    if meta:
        lines = meta if isinstance(meta, list) else [meta]
        if len(lines) > 2:
            raise _err(ctx, "`meta` takes at most two short lines")
        variants: list[tuple[str, bool]] = []
        for i, m in enumerate(lines):
            variants.extend(_variants(ctx, m, f"meta[{i}]"))
        _tx(ctx, MX + CW - meta_w, y - 4, meta_w, 58, variants, size=17, color=B["faint"], lh=1.5, align="right", anchor="bottom", name="text:meta")
    add_rect(ctx.slide, MX, HEAD_RULE_Y, CW, 1, fill=B["line"], name="deco:head-rule")
    intro = ctx.copy(s, "intro")
    if intro:
        ih = _est(intro, 20, CW, 1.5)
        if ih > 20 * 1.5 * 3 + 1:
            raise _err(ctx, "`intro` is at most three lines; move the rest to the body or the notes")
        _tx(ctx, MX, BODY_TOP, CW, ih, intro, size=20, color=B["muted"], lh=1.5, name="text:intro")
        return BODY_TOP + ih + 24
    return BODY_TOP


def _chrome(ctx: Ctx) -> None:
    """Progress bar (top), compact lockup + label (bottom-left), page counter (bottom-right)."""
    if ctx.total:
        add_rect(ctx.slide, 0, 0, W * ctx.page_number / ctx.total, 4, fill=B["accent"], name="deco:progress")
    y = FOOTER_TOP
    lw = 0.0
    if ctx.show_logo:
        pic = add_logo(ctx.slide, "small", theme.FOOTER["left"], y, theme.FOOTER["logo_h"], name="logo:small")
        lw = pic.width / theme.EMU_PER_PX + theme.FOOTER["gap"]
    label = ctx.footer_label
    if label and ctx.lang == "bi":
        label = label[:1]
    if label:
        _tx(ctx, theme.FOOTER["left"] + lw, y, 1000, theme.FOOTER["logo_h"], label, size=SMALL, color=B["faint"], caps=True, ls=0.9, anchor="middle", lh=1.0, name="footer:label")
    counter = f"{ctx.page_number:02d} / {ctx.total:02d}" if ctx.total else f"{ctx.page_number:02d}"
    add_text(ctx.slide, W - theme.FOOTER["right"] - 200, y, 200, theme.FOOTER["logo_h"], [[{"text": counter}]], size=SMALL, color=B["faint"], align="right", anchor="middle", line_height=1.0, letter_spacing=0.6, name="footer:page")


def _white(ctx: Ctx) -> None:
    from pptx.dml.color import RGBColor

    ctx.slide.background.fill.solid()
    ctx.slide.background.fill.fore_color.rgb = RGBColor.from_string(B["white"])


# ---------------------------------------------------------------------------
# blocks (used by cards, tables' bands and panels)
# ---------------------------------------------------------------------------


def _callout(ctx: Ctx, x, y, w, c: dict) -> float:
    """Tinted band with a left bar: [BADGE] title / text. Returns the bottom y."""
    if not isinstance(c, dict) or "title" not in c:
        raise _err(ctx, "`callout` needs at least a `title`")
    tone = _tone(ctx, c.get("tone"))
    title = _variants(ctx, c["title"], "callout.title")
    text = _variants(ctx, c["text"], "callout.text") if c.get("text") else []
    badge = c.get("badge")
    pad_x, pad_y = 31, 25
    bx = x + pad_x
    badge_w = text_width_em(str(badge).upper()) * 16 * 1.15 + 30 if badge else 0
    tx = bx + (badge_w + 28 if badge else 0)
    text_w = x + w - pad_x - tx
    title_h = max(36, _est(title, 30, text_w, 1.2, 600))
    th = title_h + ((8 + _est(text, 20, text_w, 1.5)) if text else 0)
    h = th + 2 * pad_y
    add_rect(ctx.slide, x, y, w, h, fill=tone["bg"], radius=13, line=tone["fg"], line_px=1, name="deco:callout")
    add_rect(ctx.slide, x, y + 8, 5, h - 16, fill=tone["fg"], radius=2, name="deco:callout-bar")
    if badge:
        _badge(ctx, bx, y + pad_y + 2, str(badge), tone)
    _tx(ctx, tx, y + pad_y, text_w, title_h, title, size=30, weight=600, color=B["ink"], lh=1.2, ls=-0.3, name="text:callout-title")
    if text:
        _tx(ctx, tx, y + pad_y + title_h + 8, text_w, th - title_h - 8, text, size=20, color=tone["text"], lh=1.5, name="text:callout")
    return y + h


def _card_est(ctx: Ctx, w: float, c: dict) -> float:
    inner = w - 50
    h = 22 + 24 + 14
    if c.get("stat") is not None:
        h += 76 + 6 + 22 + 10
    if c.get("items"):
        h += sum(_est(it, BODY, inner - BODY * 1.1, 1.6) + 6 for it in _items(ctx, c["items"], "card.items")) - 6
    if c.get("text"):
        h += _est(_variants(ctx, c["text"], "card.text"), BODY, inner, 1.6)
    return h * 1.04 + 22 + 6  # renderers wrap a little earlier than the estimate on long lists


def _card(ctx: Ctx, x, y, w, h, c: dict, idx: int) -> None:
    if not isinstance(c, dict) or "heading" not in c:
        raise _err(ctx, f"cards[{idx}] needs a `heading`")
    add_rect(ctx.slide, x, y, w, h, fill=B["white"], radius=12, line=B["line"], line_px=1, name="deco:card")
    ix, iy, inner = x + 25, y + 22, w - 50
    _kicker(ctx, ix, iy, inner, _variants(ctx, c["heading"], f"cards[{idx}].heading"), name="text:eyebrow")
    iy += 24 + 14
    if c.get("stat") is not None:
        add_text(ctx.slide, ix, iy, inner, 76, [[{"text": str(c["stat"])}]], size=64, weight=600, color=B["accent"], line_height=1.1, letter_spacing=-2, name=f"text:card{idx + 1}-stat")
        iy += 76 + 6
        if c.get("label"):
            _tx(ctx, ix, iy, inner, 22, _variants(ctx, c["label"], f"cards[{idx}].label"), size=SMALL, weight=600, color=B["faint"], caps=True, ls=1.2, lh=1.2, name=f"text:card{idx + 1}-label")
            iy += 22 + 10
    if c.get("items"):
        iy += _list_block(ctx, ix, iy, inner, _items(ctx, c["items"], f"cards[{idx}].items"), name=f"text:card{idx + 1}-items")
    if c.get("text"):
        _tx(ctx, ix, iy, inner, _est(_variants(ctx, c["text"], f"cards[{idx}].text"), BODY, inner, 1.6), _variants(ctx, c["text"], f"cards[{idx}].text"), size=BODY, color=B["muted"], lh=1.6, name=f"text:card{idx + 1}-text")


def _flow(ctx: Ctx, x, y, w, f: dict) -> float:
    """Numbered step strip with a label cell ("Identification, every case"). Returns bottom y."""
    steps = f.get("steps")
    if not steps or not isinstance(steps, list):
        raise _err(ctx, "`flow.steps` (list of {text, sub?, cond?}) is required")
    h = 78
    add_rect(ctx.slide, x, y, w, h, fill=B["white"], radius=12, line=B["line"], line_px=1, name="deco:flow")
    add_rect(ctx.slide, x, y + 8, 5, h - 16, fill=B["accent"], radius=2, name="deco:flow-bar")
    sx = x + 20
    if f.get("label"):
        _tx(ctx, x + 24, y, 176, h, _variants(ctx, f["label"], "flow.label"), size=SMALL, weight=700, color=B["accent"], caps=True, ls=1.5, lh=1.35, anchor="middle", name="text:flow-label")
        sx = x + 220
    n = len(steps)
    sw = (x + w - 20 - sx) / n
    for i, st in enumerate(steps):
        if not isinstance(st, dict) or "text" not in st:
            raise _err(ctx, f"flow.steps[{i}] needs `text`")
        cx = sx + i * sw
        _step_circle(ctx, cx + 16, y + 25, i + 1, cond=bool(st.get("cond")))
        _tx(ctx, cx + 58, y + 16, sw - 80, 26, _variants(ctx, st["text"], f"flow.steps[{i}].text"), size=18, color=B["ink"], lh=1.3, name=f"text:step{i + 1}")
        if st.get("sub"):
            _tx(ctx, cx + 58, y + 42, sw - 80, 22, _variants(ctx, st["sub"], f"flow.steps[{i}].sub"), size=SMALL, color=B["faint"], lh=1.3, name=f"text:step{i + 1}-sub")
        if i < n - 1:
            add_text(ctx.slide, cx + sw - 18, y + 22, 16, 34, [[{"text": "›"}]], size=24, color=B["line"], align="center", anchor="middle", line_height=1.0, name="deco:arrow")
    return y + h


def _rail(ctx: Ctx, x, y, w, r: dict) -> float:
    """Date rail: LABEL  [key note] [key note] ...  note-right. Returns bottom y."""
    items = r.get("items")
    if not items or not isinstance(items, list):
        raise _err(ctx, "`rail.items` (list of {key, note?, tone?}) is required")
    h = 36
    cx = x
    if r.get("label"):
        label = _variants(ctx, r["label"], "rail.label")
        lw = _width(label, SMALL, 700) * 1.5 + len(joined_plain(label)) * 1.8 + 16
        _tx(ctx, x, y, lw, h, label, size=SMALL, weight=700, color=B["faint"], caps=True, ls=1.8, lh=1.0, anchor="middle", name="text:rail-label")
        cx = x + lw + 12
    for i, it in enumerate(items):
        if not isinstance(it, dict) or "key" not in it:
            raise _err(ctx, f"rail.items[{i}] needs `key`")
        key = joined_plain(_variants(ctx, it["key"], f"rail.items[{i}].key"))
        note = joined_plain(_variants(ctx, it["note"], f"rail.items[{i}].note")) if it.get("note") else ""
        tone = _tone(ctx, it["tone"]) if it.get("tone") else None
        pw = text_width_em(key) * 17 * 1.06 + (text_width_em(note) * 15.5 + 10 if note else 0) + 36
        shp = add_rect(ctx.slide, cx, y, pw, h, fill=tone["bg"] if tone else B["white"], radius=h / 2, line=tone["fg"] if tone else B["line"], line_px=1, name="pill:rail")
        runs = [(key, 17, 600, tone["fg"] if tone else B["ink"])]
        if note:
            runs.append(("  " + note, 15.5, 400, tone["fg"] if tone else B["faint"]))
        _shape_runs(shp, runs, align="center", inset_l=16, inset_r=16)
        cx += pw + 10
    if r.get("note"):
        _tx(ctx, cx + 10, y, x + w - cx - 10, h, _variants(ctx, r["note"], "rail.note"), size=16, color=B["faint"], lh=1.0, align="right", anchor="middle", name="text:rail-note")
    return y + h


_CELL_KINDS = ("name", "text", "muted", "date", "status")


def _cell_lines(ctx: Ctx, cell, kind: str, colw: float) -> tuple[float, list, list]:
    """Return (height, main_variants, sub_variants) for a cell."""
    sub: list = []
    if isinstance(cell, dict) and ("text" in cell or "sub" in cell or "label" in cell):
        main = _variants(ctx, cell.get("text", cell.get("label", "")), "table.cell") if (cell.get("text") or cell.get("label")) else []
        sub = _items(ctx, cell.get("sub"), "table.cell.sub")
    else:
        main = _variants(ctx, cell, "table.cell") if cell not in (None, "") else []
    weight = 600 if kind in ("name", "date", "status") else 400
    size = 17 if kind == "muted" else 18
    inner = colw - 32 - (20 if kind == "name" else 0)
    h = _est(main, size, inner, 1.4, weight) if main else size * 1.4
    if sub:
        h += 6 + sum(_est(s, 15, inner - 12, 1.5) for s in sub)
    return h, main, sub


def _table(ctx: Ctx, x, y, w, t: dict) -> float:
    """Reference table: uppercase faint headers, hairline rows, status dot in the name
    column, hot dates in stop red, status pills as tinted text. Returns the height."""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    cols = t.get("columns")
    rows = t.get("rows")
    if not cols or not isinstance(cols, list):
        raise _err(ctx, "`table.columns` (list of {header, width, kind}) is required")
    if not rows or not isinstance(rows, list):
        raise _err(ctx, "`table.rows` is required")
    widths = []
    for i, c in enumerate(cols):
        if not isinstance(c, dict) or "header" not in c or "width" not in c:
            raise _err(ctx, f"table.columns[{i}] needs `header` and `width` (fraction of the table width)")
        if c.get("kind", "text") not in _CELL_KINDS:
            raise _err(ctx, f"table.columns[{i}].kind must be one of {_CELL_KINDS}")
        widths.append(float(c["width"]))
    if abs(sum(widths) - 1.0) > 0.02:
        raise _err(ctx, f"table column widths must sum to 1.0 (got {sum(widths):.2f})")
    kinds = [c.get("kind", "text") for c in cols]
    px_w = [w * f for f in widths]
    # measure
    header_h = 40
    row_meta = []
    for ri, r in enumerate(rows):
        cells = r.get("cells") if isinstance(r, dict) else r
        if not isinstance(cells, list) or len(cells) != len(cols):
            raise _err(ctx, f"table.rows[{ri}] needs {len(cols)} cells")
        status = r.get("status") if isinstance(r, dict) else None
        if status is not None and status not in TONES:
            raise _err(ctx, f"table.rows[{ri}].status must be one of {sorted(TONES)}")
        measured = [_cell_lines(ctx, c, kinds[ci], px_w[ci]) for ci, c in enumerate(cells)]
        rh = max(m[0] for m in measured) + 24
        row_meta.append((status, measured, rh))
    total_h = header_h + sum(m[2] for m in row_meta)
    if y + total_h > BODY_BOTTOM:
        raise _err(ctx, f"the table needs ~{total_h:.0f}px but only {BODY_BOTTOM - y:.0f}px remain; a briefing slide holds about 8 short rows - drop rows, shorten cells or split the topic")

    gf = ctx.slide.shapes.add_table(len(rows) + 1, len(cols), theme.px(x), theme.px(y), theme.px(w), theme.px(total_h))
    gf.name = "table:brief"
    tbl = gf.table
    for flag in ("first_row", "horz_banding", "vert_banding", "last_row", "first_col", "last_col"):
        setattr(tbl, flag, False)
    for ci, cw in enumerate(px_w):
        tbl.columns[ci].width = theme.px(cw)
    tbl.rows[0].height = theme.px(header_h)
    for ri, (_, _, rh) in enumerate(row_meta, start=1):
        tbl.rows[ri].height = theme.px(rh)

    from .charts import _set_cell_borders

    def prep(cell, *, first: bool, top_pad: float):
        cell.fill.background()
        cell.margin_left = theme.px(13 if first else 16)
        cell.margin_right = theme.px(16)
        cell.margin_top = theme.px(top_pad)
        cell.margin_bottom = theme.px(12)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        tf = cell.text_frame
        tf.word_wrap = True
        return tf

    for ci, c in enumerate(cols):
        tf = prep(tbl.cell(0, ci), first=ci == 0, top_pad=0)
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = joined_plain(_variants(ctx, c["header"], f"table.columns[{ci}].header"))
        _apply_run_style(run, size_px=14, weight=700, color=B["faint"], alpha=None, italic=False, letter_spacing=1.5, caps=True)
        _set_cell_borders(tbl.cell(0, ci), bottom_color=B["line"], bottom_alpha=None, bottom_px=1)

    for ri, (status, measured, rh) in enumerate(row_meta, start=1):
        last = ri == len(rows)
        for ci, (h_, main, sub) in enumerate(measured):
            kind = kinds[ci]
            tf = prep(tbl.cell(ri, ci), first=ci == 0, top_pad=12)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = Pt(theme.pt(18) * 1.4)
            cell_spec = rows[ri - 1]["cells"][ci] if isinstance(rows[ri - 1], dict) else rows[ri - 1][ci]
            tone = None
            if kind == "status":
                tname = cell_spec.get("tone") if isinstance(cell_spec, dict) else None
                tone = _tone(ctx, tname) if tname else None
            if kind == "name" and status:
                r = p.add_run()
                r.text = "● "
                _apply_run_style(r, size_px=13, weight=400, color=TONES[status]["fg"], alpha=None, italic=False, letter_spacing=0, caps=False)
            if kind == "status" and tone:
                r = p.add_run()
                r.text = "● "
                _apply_run_style(r, size_px=13, weight=400, color=tone["fg"], alpha=None, italic=False, letter_spacing=0, caps=False)
            if kind == "name":
                color, weight, size = B["ink"], 600, 18
            elif kind == "date":
                hot = status == "stop"
                color, weight, size = (TONES["stop"]["fg"] if hot else B["ink"]), 600, 18
            elif kind == "status":
                color, weight, size = (tone["fg"] if tone else B["muted"]), 600, 17
            elif kind == "muted":
                color, weight, size = B["muted"], 400, 17
            else:
                color, weight, size = B["muted"], 400, 18
            if main:
                for para_i, spans in enumerate(_paras(ctx, main, color=color, weight=weight)):
                    if para_i:
                        p = tf.add_paragraph()
                        p.line_spacing = Pt(theme.pt(size) * 1.4)
                    for span in spans:
                        if span.get("br"):
                            p.add_line_break()
                            continue
                        r = p.add_run()
                        r.text = span["text"]
                        _apply_run_style(r, size_px=span.get("size", size), weight=span.get("weight", weight), color=span.get("color", color), alpha=span.get("alpha"), italic=span.get("italic", False), letter_spacing=0, caps=False, lang=span.get("lang"))
            elif not (kind in ("name", "status") and (status or tone)):
                r = p.add_run()
                r.text = "–"
                _apply_run_style(r, size_px=size, weight=400, color=B["faint"], alpha=None, italic=False, letter_spacing=0, caps=False)
            for s in sub:
                sp = tf.add_paragraph()
                sp.line_spacing = Pt(theme.pt(15) * 1.5)
                sp.space_before = Pt(theme.pt(4))
                for spans in _paras(ctx, s, color=B["faint"]):
                    for span in spans:
                        if span.get("br"):
                            sp.add_line_break()
                            continue
                        r = sp.add_run()
                        r.text = span["text"]
                        _apply_run_style(r, size_px=15, weight=span.get("weight", 400), color=span.get("color", B["faint"]), alpha=None, italic=False, letter_spacing=0, caps=False, lang=span.get("lang"))
            _set_cell_borders(tbl.cell(ri, ci), bottom_color=None if last else B["line_soft"], bottom_alpha=None, bottom_px=0 if last else 1)
    return total_h


def _band_box_est(ctx: Ctx, w: float, b: dict) -> float:
    inner = w - 2 * 21 - 6
    h = 20 + 6
    if b.get("asks"):
        h += sum(_est(it, 18, inner - 24, 1.5) + 3 for it in _items(ctx, b["asks"], "band.asks")) - 3
    if b.get("text"):
        h += _est(_variants(ctx, b["text"], "band.text"), 18, inner, 1.45)
    return h + 2 * 15


def _band_box(ctx: Ctx, x, y, w, h, b: dict, idx: int) -> None:
    if not isinstance(b, dict) or "key" not in b:
        raise _err(ctx, f"band[{idx}] needs a `key` label")
    tinted = bool(b.get("asks")) or bool(b.get("tinted"))
    add_rect(ctx.slide, x, y, w, h, fill=B["accent_xlt"] if tinted else B["white"], radius=10, line=B["line"], line_px=1, name="deco:band")
    add_rect(ctx.slide, x, y + 6, 4, h - 12, fill=B["accent"], radius=2, name="deco:band-bar")
    ix, iy, inner = x + 21 + 6, y + 15, w - 2 * 21 - 6
    _tx(ctx, ix, iy, inner, 20, _variants(ctx, b["key"], f"band[{idx}].key"), size=SMALL, weight=700, color=B["accent"], caps=True, ls=1.8, lh=1.2, name=f"text:band{idx + 1}-key")
    iy += 26
    if b.get("asks"):
        _list_block(ctx, ix, iy, inner, _items(ctx, b["asks"], f"band[{idx}].asks"), size=18, lh=1.5, numbered=True, gap=3, name=f"text:band{idx + 1}-asks")
    if b.get("text"):
        v = _variants(ctx, b["text"], f"band[{idx}].text")
        _tx(ctx, ix, iy, inner, _est(v, 18, inner, 1.45), v, size=18, color=B["muted"], lh=1.45, name=f"text:band{idx + 1}-text")


def _band(ctx: Ctx, x, y, w, boxes: list) -> float:
    if not isinstance(boxes, list) or not 1 <= len(boxes) <= 2:
        raise _err(ctx, "`band` takes one or two boxes ({key, text} or {key, asks: [...]})")
    gap = 20
    if len(boxes) == 2:
        widths = [(w - gap) / 2.35, (w - gap) * 1.35 / 2.35]
    else:
        widths = [w]
    heights = [_band_box_est(ctx, bw, b) for bw, b in zip(widths, boxes)]
    h = max(heights)
    if y + h > BODY_BOTTOM:
        raise _err(ctx, f"the bottom band needs ~{h:.0f}px but only {BODY_BOTTOM - y:.0f}px remain below the table; shorten the table or the band")
    cx = x
    for i, (bw, b) in enumerate(zip(widths, boxes)):
        _band_box(ctx, cx, y, bw, h, b, i)
        cx += bw + gap
    return y + h


def _timeline(ctx: Ctx, x, y, w, t: dict) -> float:
    """Phase cards in a row (the "delivery plan" strip): kicker label, then per phase a
    bordered card with title, text and a date pill. Returns the height."""
    phases = t.get("phases")
    if not phases or not isinstance(phases, list) or not 2 <= len(phases) <= 4:
        raise _err(ctx, "`timeline.phases` takes two to four {title, text?, date?, tone?} entries")
    h_label = 0.0
    if t.get("label"):
        _kicker(ctx, x, y, w, _variants(ctx, t["label"], "timeline.label"), name="text:eyebrow", size=SMALL)
        h_label = 24 + 12
    n = len(phases)
    gap = 20
    pw = (w - (n - 1) * gap) / n
    inner = pw - 44
    heights = []
    for i, ph in enumerate(phases):
        if not isinstance(ph, dict) or "title" not in ph:
            raise _err(ctx, f"timeline.phases[{i}] needs a `title`")
        th = _est(_variants(ctx, ph["title"], f"timeline.phases[{i}].title"), 18, inner - (150 if ph.get("date") else 0), 1.3, 600)
        tx_h = _est(_variants(ctx, ph["text"], f"timeline.phases[{i}].text"), SMALL, inner, 1.5) if ph.get("text") else 0
        heights.append(18 + th + (6 + tx_h if tx_h else 0) + 18)
    ch = max(heights)
    for i, ph in enumerate(phases):
        cx = x + i * (pw + gap)
        cy = y + h_label
        add_rect(ctx.slide, cx, cy, pw, ch, fill=B["white"], radius=10, line=B["line"], line_px=1, name="deco:phase")
        add_rect(ctx.slide, cx, cy + 6, 4, ch - 12, fill=B["accent"], radius=2, name="deco:phase-bar")
        title = _variants(ctx, ph["title"], f"timeline.phases[{i}].title")
        date = joined_plain(_variants(ctx, ph["date"], f"timeline.phases[{i}].date")) if ph.get("date") else ""
        th = _est(title, 18, inner - (150 if date else 0), 1.3, 600)
        _tx(ctx, cx + 22, cy + 18, inner - (150 if date else 0), th, title, size=18, weight=600, color=B["ink"], lh=1.3, name=f"text:phase{i + 1}-title")
        if date:
            tone = _tone(ctx, ph.get("tone"))
            pw_ = text_width_em(date) * 14.5 * 1.2 + 47
            _pill(ctx, cx + pw - 16 - pw_, cy + 16, date, tone, h=26, size=14.5)
        if ph.get("text"):
            v = _variants(ctx, ph["text"], f"timeline.phases[{i}].text")
            _tx(ctx, cx + 22, cy + 18 + th + 6, inner, _est(v, SMALL, inner, 1.5), v, size=SMALL, color=B["muted"], lh=1.5, name=f"text:phase{i + 1}-text")
    return h_label + ch


# --- panel blocks -----------------------------------------------------------


def _block_matrix(ctx: Ctx, x, y, w, b: dict) -> float:
    rows = b.get("rows")
    if not rows or not isinstance(rows, list):
        raise _err(ctx, "matrix block needs `rows` (list of {name, status: {label, tone}})")
    header = b.get("header") or ["Agent", "Status"]
    hh, rh = 38, 42
    h = hh + rh * len(rows)
    add_rect(ctx.slide, x, y, w, h, fill=B["white"], radius=10, line=B["line"], line_px=1, name="deco:matrix")
    _tx(ctx, x + 21, y, w / 2 - 21, hh, _variants(ctx, header[0], "matrix.header[0]"), size=14, weight=700, color=B["faint"], caps=True, ls=1.5, lh=1.0, anchor="middle", name="text:matrix-h1")
    _tx(ctx, x + w / 2, y, w / 2 - 21, hh, _variants(ctx, header[1], "matrix.header[1]"), size=14, weight=700, color=B["faint"], caps=True, ls=1.5, lh=1.0, align="right", anchor="middle", name="text:matrix-h2")
    add_rect(ctx.slide, x + 1, y + hh, w - 2, 1, fill=B["line"], name="deco:matrix-rule")
    ry = y + hh
    for i, r in enumerate(rows):
        if not isinstance(r, dict) or "name" not in r or "status" not in r:
            raise _err(ctx, f"matrix.rows[{i}] needs `name` and `status: {{label, tone}}`")
        if i % 2 == 1:
            add_rect(ctx.slide, x + 1, ry, w - 2, rh, fill=B["accent_xlt"], name="deco:matrix-band")
        _tx(ctx, x + 21, ry, w / 2 + 40, rh, _variants(ctx, r["name"], f"matrix.rows[{i}].name"), size=18, color=B["ink"], lh=1.0, anchor="middle", name=f"text:matrix-r{i + 1}")
        st = r["status"]
        if not isinstance(st, dict) or "label" not in st or "tone" not in st:
            raise _err(ctx, f"matrix.rows[{i}].status needs `label` and `tone`")
        label = joined_plain(_variants(ctx, st["label"], f"matrix.rows[{i}].status.label"))
        pw = text_width_em(label) * 15.5 * 1.2 + 47
        _pill(ctx, x + w - 14 - pw, ry + (rh - 28) / 2, label, _tone(ctx, st["tone"]), h=28)
        if i < len(rows) - 1:
            add_rect(ctx.slide, x + 1, ry + rh - 1, w - 2, 1, fill=B["line_soft"], name="deco:matrix-hair")
        ry += rh
    return h


def _block_box(ctx: Ctx, x, y, w, b: dict) -> float:
    if "title" not in b:
        raise _err(ctx, "box block needs a `title`")
    title = _variants(ctx, b["title"], "box.title")
    items = _items(ctx, b.get("items"), "box.items")
    text = _variants(ctx, b["text"], "box.text") if b.get("text") else []
    inner = w - 50
    tw = _width(title, 21, 600) * 1.3
    pill_w = 0.0
    if b.get("pill"):
        p = b["pill"]
        if not isinstance(p, dict) or "label" not in p or "tone" not in p:
            raise _err(ctx, "box.pill needs `label` and `tone`")
        pill_w = text_width_em(joined_plain(_variants(ctx, p["label"], "box.pill.label"))) * 14.5 * 1.2 + 47
    pill_below = pill_w and (tw + 14 + pill_w > inner)
    title_h = 28 + (34 if pill_below else 0)
    body_h = (sum(_est(it, BODY, inner - BODY * 1.1, 1.6) + 6 for it in items) - 6) if items else 0
    if text:
        body_h += (6 if items else 0) + _est(text, BODY, inner, 1.6)
    h = 22 + title_h + (12 if (items or text) else 0) + body_h + 22
    add_rect(ctx.slide, x, y, w, h, fill=B["white"], radius=12, line=B["line"], line_px=1, name="deco:box")
    ix, iy = x + 25, y + 22
    _tx(ctx, ix, iy, inner, 28, title, size=21, weight=600, color=B["ink"], lh=1.3, anchor="middle", name="text:box-title")
    if pill_w:
        p = b["pill"]
        px_, py_ = (ix, iy + 32) if pill_below else (ix + tw + 14, iy + 1)
        _pill(ctx, px_, py_, joined_plain(_variants(ctx, p["label"], "box.pill.label")), _tone(ctx, p["tone"]), h=26, size=14.5)
    iy += title_h + 12
    if items:
        iy += _list_block(ctx, ix, iy, inner, items, name="text:box-items") + 6
    if text:
        _tx(ctx, ix, iy, inner, _est(text, BODY, inner, 1.6), text, size=BODY, color=B["muted"], lh=1.6, name="text:box-text")
    return h


def _block_ask(ctx: Ctx, x, y, w, b: dict) -> float:
    if "title" not in b or not b.get("items"):
        raise _err(ctx, "ask block needs `title` and `items`")
    items = _items(ctx, b["items"], "ask.items")
    inner = w - 50
    body_h = sum(_est(it, 18, inner - 24, 1.55) + 5 for it in items) - 5
    h = 21 + 20 + 12 + body_h + 21
    add_rect(ctx.slide, x, y, w, h, fill=B["accent_xlt"], radius=12, line=B["accent"], line_px=1, name="deco:ask")
    ix, iy = x + 25, y + 21
    _tx(ctx, ix, iy, inner, 20, _variants(ctx, b["title"], "ask.title"), size=15.5, weight=700, color=B["accent"], caps=True, ls=1.7, lh=1.2, name="text:ask-title")
    iy += 20 + 12
    _list_block(ctx, ix, iy, inner, items, size=18, lh=1.55, numbered=True, gap=5, name="text:ask-items")
    return h


def _block_note(ctx: Ctx, x, y, w, b: dict) -> float:
    if "text" not in b:
        raise _err(ctx, "note block needs `text`")
    v = _variants(ctx, b["text"], "note.text")
    h = _est(v, 16, w, 1.5)
    _tx(ctx, x, y, w, h, v, size=16, color=B["faint"], lh=1.5, name="text:note")
    return h


def _block_list(ctx: Ctx, x, y, w, b: dict) -> float:
    items = _items(ctx, b.get("items"), "list.items")
    if not items:
        raise _err(ctx, "list block needs `items`")
    return _list_block(ctx, x, y, w, items, name="text:list")


def _block_text(ctx: Ctx, x, y, w, b: dict) -> float:
    if "text" not in b:
        raise _err(ctx, "text block needs `text`")
    v = _variants(ctx, b["text"], "text.text")
    h = _est(v, BODY, w, 1.6)
    _tx(ctx, x, y, w, h, v, size=BODY, color=B["muted"], lh=1.6, name="text:body")
    return h


def _block_steps(ctx: Ctx, x, y, w, b: dict) -> float:
    return _flow(ctx, x, y, w, b) - y


def _block_table(ctx: Ctx, x, y, w, b: dict) -> float:
    return _table(ctx, x, y, w, b)


def _block_timeline(ctx: Ctx, x, y, w, b: dict) -> float:
    return _timeline(ctx, x, y, w, b)


_BLOCKS = {
    "table": _block_table,
    "timeline": _block_timeline,
    "matrix": _block_matrix,
    "box": _block_box,
    "ask": _block_ask,
    "note": _block_note,
    "list": _block_list,
    "text": _block_text,
    "steps": _block_steps,
}


def _block(ctx: Ctx, x, y, w, b: dict) -> float:
    if not isinstance(b, dict) or "type" not in b:
        raise _err(ctx, f"each block needs a `type` (one of {sorted(_BLOCKS)})")
    if b["type"] not in _BLOCKS:
        raise _err(ctx, f"unknown block type {b['type']!r}; use one of {sorted(_BLOCKS)}")
    return _BLOCKS[b["type"]](ctx, x, y, w, b)


# ---------------------------------------------------------------------------
# templates
# ---------------------------------------------------------------------------


def brief_cover(ctx: Ctx, s: dict) -> None:
    """Eyebrow, short accent rule, big SemiBold title (accent span via ~), subtitle,
    up to three section cards, concentric mark at the right edge, quiet foot line."""
    _white(ctx)
    cx, cy = 1545, 540
    add_oval(ctx.slide, cx - 375, cy - 375, 750, 750, fill=None, line=B["line"], line_px=1.5, name="deco:mark-outer")
    add_oval(ctx.slide, cx - 267, cy - 267, 534, 534, fill=None, line=B["line_soft"], line_px=1.5, name="deco:mark-inner")
    add_oval(ctx.slide, cx - 159, cy - 159, 318, 318, fill=B["accent_xlt"], name="deco:mark-disc")
    X, XW = 114, 1400
    if ctx.show_logo:
        add_logo(ctx.slide, "primary", X, TOP, 44)
    y = 250
    eyebrow = ctx.copy(s, "eyebrow")
    if eyebrow:
        _kicker(ctx, X, y, XW, eyebrow, name="text:eyebrow", size=16)
    y += 28 + 14
    add_rect(ctx.slide, X, y, 81, 4, fill=B["accent"], name="deco:rule")
    y += 4 + 36
    title = ctx.copy(s, "title", required=True)
    th = _est(title, 84, XW, 1.03, 600)
    if th > 84 * 1.03 * 2 + 1:
        ctx.warn("cover title exceeds two lines at 84px; shorten it")
    _tx(ctx, X, y, XW, th, title, size=84, weight=600, color=B["ink"], lh=1.03, ls=-2.5, name="text:title")
    y += th + 22
    sub = ctx.copy(s, "subtitle")
    if sub:
        sh = _est(sub, 28, XW, 1.3)
        _tx(ctx, X, y, XW, sh, sub, size=28, color=B["muted"], lh=1.3, name="text:subtitle")
        y += sh
    sections = s.get("sections") or []
    if not isinstance(sections, list) or len(sections) > 3:
        raise _err(ctx, "`sections` takes up to three {heading, text} cards")
    if sections:
        y += 60
        n = len(sections)
        gap = 21
        sw = (W - 2 * X - (n - 1) * gap) / n
        inner = sw - 54
        texts = [_variants(ctx, sec.get("text", ""), f"sections[{i}].text") if sec.get("text") else [] for i, sec in enumerate(sections)]
        ch = max(22 + 30 + 6 + (_est(t, 17, inner, 1.5) if t else 0) + 24 for t in texts)
        for i, sec in enumerate(sections):
            if not isinstance(sec, dict) or "heading" not in sec:
                raise _err(ctx, f"sections[{i}] needs a `heading`")
            sx = X + i * (sw + gap)
            add_rect(ctx.slide, sx, y, sw, ch, fill=B["white"], radius=13, line=B["line"], line_px=1, name="deco:section-card")
            add_rect(ctx.slide, sx + 12, y, sw - 24, 4, fill=B["accent"], radius=2, name="deco:section-top")
            _tx(ctx, sx + 27, y + 22, inner, 30, _variants(ctx, sec["heading"], f"sections[{i}].heading"), size=22, weight=600, color=B["ink"], lh=1.3, name=f"text:section{i + 1}-heading")
            if texts[i]:
                _tx(ctx, sx + 27, y + 22 + 30 + 6, inner, ch - 82, texts[i], size=17, color=B["muted"], lh=1.5, name=f"text:section{i + 1}-text")
        y += ch
    if y > 960:
        raise _err(ctx, f"the cover content flows to y={y:.0f}px and would reach the foot line; shorten the title, subtitle or section cards")
    fy = H - 63 - 24
    if s.get("foot_left"):
        _tx(ctx, X, fy, 820, 24, ctx.copy(s, "foot_left"), size=16, color=B["faint"], caps=True, ls=0.7, lh=1.2, anchor="middle", name="text:meta")
    if s.get("foot_right"):
        _tx(ctx, W - X - 820, fy, 820, 24, ctx.copy(s, "foot_right"), size=16, color=B["faint"], caps=True, ls=0.7, lh=1.2, align="right", anchor="middle", name="text:meta-right")


def brief_cards(ctx: Ctx, s: dict) -> None:
    """Head, optional callout band, 2-4 cards (kicker + bullets / text / stat), optional note."""
    _white(ctx)
    y = _head(ctx, s)
    if s.get("callout"):
        y = _callout(ctx, MX, y, CW, s["callout"]) + 32
    cards = s.get("cards")
    if not cards or not isinstance(cards, list) or not 1 <= len(cards) <= 12:
        raise _err(ctx, "`cards` takes one to twelve {heading, items|text|stat} cards (one row up to four, then rows of three or four)")
    note = ctx.copy(s, "note")
    note_h = (_est(note, BODY, CW, 1.6) + 20) if note else 0
    bottom = BODY_BOTTOM - (note_h + 24 if note else 0)
    timeline = s.get("timeline")
    n = len(cards)
    gap = 22
    per_row = n if n <= 4 else (3 if n in (5, 6, 9) else 4)
    cw = (CW - (per_row - 1) * gap) / per_row
    rows = [cards[i : i + per_row] for i in range(0, n, per_row)]
    row_h = [max(max(_card_est(ctx, cw, c) for c in r), 200 if len(rows) == 1 else 120) for r in rows]
    cards_h = sum(row_h) + (len(rows) - 1) * gap
    if timeline:
        # measure the strip on a scratch estimate: cards must leave room for it
        tl_h = (24 + 12 if timeline.get("label") else 0) + 100
    else:
        tl_h = 0
    if y + cards_h + (32 + tl_h if timeline else 0) > bottom:
        raise _err(ctx, f"the cards need ~{cards_h:.0f}px but only {bottom - y - (32 + tl_h if timeline else 0):.0f}px remain; trim the bullets (a briefing card holds 3-5 short lines), drop the callout or split the slide")
    cy = y
    idx = 0
    for r, rh in zip(rows, row_h):
        for k, c in enumerate(r):
            _card(ctx, MX + k * (cw + gap), cy, cw, rh, c, idx)
            idx += 1
        cy += rh + gap
    if timeline:
        th = _timeline(ctx, MX, cy - gap + 32, CW, timeline)
        if cy - gap + 32 + th > bottom:
            raise _err(ctx, "the timeline strip does not fit under the cards; shorten the phase texts or the cards")
    if note:
        ny = BODY_BOTTOM - note_h
        add_rect(ctx.slide, MX, ny, CW, 1, fill=B["line"], name="deco:note-rule")
        _tx(ctx, MX, ny + 20, CW, note_h - 20, note, size=BODY, color=B["muted"], lh=1.6, name="text:footnote")
    _chrome(ctx)


def brief_table(ctx: Ctx, s: dict) -> None:
    """Head, optional step strip, optional date rail, the table, optional bottom band."""
    _white(ctx)
    y = _head(ctx, s)
    if s.get("flow"):
        y = _flow(ctx, MX, y, CW, s["flow"]) + 20
    if s.get("rail"):
        y = _rail(ctx, MX, y, CW, s["rail"]) + 18
    table = s.get("table")
    if not table:
        raise _err(ctx, "`table: {columns, rows}` is required")
    y += _table(ctx, MX, y, CW, table)
    if s.get("band"):
        _band(ctx, MX, y + 26, CW, s["band"])
    _chrome(ctx)


def brief_panels(ctx: Ctx, s: dict) -> None:
    """Head + one to three panels, each a kicker title and a stack of blocks
    (matrix / box / ask / note / list / text / steps)."""
    _white(ctx)
    y = _head(ctx, s)
    panels = s.get("panels")
    if not panels or not isinstance(panels, list) or not 1 <= len(panels) <= 3:
        raise _err(ctx, "`panels` takes one to three {title, blocks: [...]} panels")
    n = len(panels)
    gap = 36
    pw = (CW - (n - 1) * gap) / n
    for i, p in enumerate(panels):
        if not isinstance(p, dict) or "title" not in p or not p.get("blocks"):
            raise _err(ctx, f"panels[{i}] needs `title` and a non-empty `blocks` list")
        x = MX + i * (pw + gap)
        py = y
        title = _variants(ctx, p["title"], f"panels[{i}].title")
        tw = _width(title, 16, 600) * 1.12 + 6
        _kicker(ctx, x, py, pw, title, name="text:eyebrow", size=16)
        add_rect(ctx.slide, x + tw + 14, py + 11, max(0, pw - tw - 14), 1, fill=B["line"], name="deco:panel-rule")
        py += 24 + 16
        for j, b in enumerate(p["blocks"]):
            if j:
                py += 18
            py += _block(ctx, x, py, pw, b)
            if py > BODY_BOTTOM:
                raise _err(ctx, f"panels[{i}] overflows at block {j + 1} (y={py:.0f}px, limit {BODY_BOTTOM}px); fewer or shorter blocks")
    _chrome(ctx)


def brief_chart(ctx: Ctx, s: dict) -> None:
    """Head, optional text column on the left, native chart in accent tints, optional note."""
    _white(ctx)
    y = _head(ctx, s)
    chart = s.get("chart")
    if not chart:
        raise _err(ctx, "`chart: {type, categories, series}` is required")
    note = ctx.copy(s, "note")
    note_h = (_est(note, 16, CW, 1.5) + 16) if note else 0
    bottom = BODY_BOTTOM - (note_h + 16 if note else 0)
    body = s.get("body")
    x, w = MX, CW
    if body:
        items = body if isinstance(body, list) else [body]
        paragraphs: list[list[dict]] = []
        est = 0.0
        for i, b in enumerate(items):
            v = _variants(ctx, b, f"body[{i}]")
            paragraphs.extend(_paras(ctx, v, color=B["muted"]))
            est += _est(v, BODY, 520, 1.6) + 16
        if y + est > bottom:
            raise _err(ctx, "the body text beside the chart is too long; keep it to two short paragraphs")
        add_text(ctx.slide, MX, y, 520, est, paragraphs, size=BODY, color=B["muted"], line_height=1.6, para_space_after=16, name="text:body")
        x, w = MX + 520 + 48, CW - 520 - 48
    add_chart(ctx.slide, x, y, w, bottom - y, chart, palette=theme.BRIEF_CHART_SERIES, axis_ink=B["faint"], grid=B["line"], label_ink=B["ink"])
    if note:
        _tx(ctx, MX, BODY_BOTTOM - note_h + 16, CW, note_h - 16, note, size=16, color=B["faint"], lh=1.5, name="text:footnote")
    _chrome(ctx)


def brief_statement(ctx: Ctx, s: dict) -> None:
    """A single statement (demo placeholder, closing page): eyebrow, big title, subtitle,
    concentric mark, primary lockup, footer chrome."""
    _white(ctx)
    cx, cy = 1545, 540
    add_oval(ctx.slide, cx - 375, cy - 375, 750, 750, fill=None, line=B["line"], line_px=1.5, name="deco:mark-outer")
    add_oval(ctx.slide, cx - 267, cy - 267, 534, 534, fill=None, line=B["line_soft"], line_px=1.5, name="deco:mark-inner")
    add_oval(ctx.slide, cx - 159, cy - 159, 318, 318, fill=B["accent_xlt"], name="deco:mark-disc")
    X, XW = 114, 1200
    y = 330
    eyebrow = ctx.copy(s, "eyebrow")
    if eyebrow:
        _kicker(ctx, X, y, XW, eyebrow, name="text:eyebrow", size=16)
        y += 28 + 14
        add_rect(ctx.slide, X, y, 81, 4, fill=B["accent"], name="deco:rule")
        y += 4 + 36
    title = ctx.copy(s, "title", required=True)
    th = _est(title, 84, XW, 1.03, 600)
    _tx(ctx, X, y, XW, th, title, size=84, weight=600, color=B["ink"], lh=1.03, ls=-2.5, name="text:title")
    y += th + 22
    sub = ctx.copy(s, "subtitle")
    if sub:
        _tx(ctx, X, y, XW, _est(sub, 28, XW, 1.3), sub, size=28, color=B["muted"], lh=1.3, name="text:subtitle")
    _chrome(ctx)


BRIEF_TEMPLATES = {
    "brief_cover": brief_cover,
    "brief_statement": brief_statement,
    "brief_cards": brief_cards,
    "brief_table": brief_table,
    "brief_panels": brief_panels,
    "brief_chart": brief_chart,
}
BRIEF_COVERS = {"brief_cover"}
