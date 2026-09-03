"""The nine NBG slide templates from `slide-templates.jsx`, plus the chart, table and
back-cover companions, rebuilt as native PowerPoint compositions.

Every coordinate below is the artboard value from the JSX (1920x1080 px). Text that
flows in HTML (title -> subtitle -> meta) is positioned from a line-count estimate so
the vertical rhythm matches the browser rendering; the validator re-checks the boxes.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from . import theme
from .charts import add_chart, add_table
from .primitives import (
    add_gradient_rect,
    add_hairline,
    add_logo,
    add_diagonal_hairlines,
    add_picture_cover,
    add_rect,
    add_shadow,
    add_text,
    estimate_height,
    estimate_lines,
)
from .text import SpecError, joined_plain, resolve_copy, spans_for

MARGIN = 90
CONTENT_W = theme.ARTBOARD_W - 2 * MARGIN  # 1740
FOOTER_TOP = theme.ARTBOARD_H - theme.FOOTER["bottom"] - theme.FOOTER["logo_h"]  # 1016
BODY_LIMIT = FOOTER_TOP - theme.FOOTER_CLEARANCE  # 944 - nothing below this except footers


@dataclass
class Ctx:
    slide: object
    lang: str
    template: str
    accent: str
    surface: str
    show_logo: bool
    page_number: int
    footer_label: object | None
    index: int
    warnings: list[str] = field(default_factory=list)
    total: int = 0  # slides in the deck (briefing chrome: progress bar, "02 / 07" counter)

    @property
    def dark(self) -> bool:
        return theme.is_dark(self.surface)

    def warn(self, msg: str) -> None:
        self.warnings.append(f"slide {self.index + 1} ({self.template}): {msg}")

    def copy(self, s: dict, key: str, *, required: bool = False):
        val = s.get(key)
        if val is None:
            if required:
                raise SpecError(f"slide {self.index + 1} ({self.template}): `{key}` is required")
            return []
        return resolve_copy(val, self.lang, path=f"slide {self.index + 1}.{key}")


# ---------------------------------------------------------------------------
# shared pieces
# ---------------------------------------------------------------------------


def _set_background(slide, hex6: str) -> None:
    from pptx.dml.color import RGBColor

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(hex6)


def _resolve_photo(ctx: Ctx, s: dict) -> tuple[Path, str]:
    key = s.get("photo")
    if key is None:
        key = theme.TEMPLATE_DEFAULT_PHOTO.get(ctx.template)
    if key is None:
        raise SpecError(f"slide {ctx.index + 1} ({ctx.template}): `photo` is required")
    if key in theme.PHOTOS:
        return theme.PHOTOS[key], f"nbg-photo:{key}"
    p = Path(str(key)).expanduser()
    if not p.is_file():
        raise SpecError(f"slide {ctx.index + 1} ({ctx.template}): photo {key!r} is neither a bundled photo {sorted(theme.PHOTOS)} nor an existing file")
    return p, f"user-photo:{p.name}"


def _text(ctx: Ctx, x, y, w, h, variants, *, size, weight=400, color, alpha=None, lh=1.2, ls=0.0, align="left", anchor="top", caps=False, accent_color=None, accent_italic=False, accent_weight=None, emphasis_weight=400, secondary_size=None, name, para_after=0):
    if not variants:
        return None
    paragraphs = spans_for(
        variants,
        base_color=color,
        accent_color=accent_color or ctx.accent,
        base_weight=weight,
        base_alpha=alpha,
        emphasis_weight=emphasis_weight,
        accent_italic=accent_italic,
        accent_weight=accent_weight,
        secondary_size=secondary_size,
        lang_hint=ctx.lang,
    )
    return add_text(ctx.slide, x, y, w, h, paragraphs, size=size, weight=weight, color=color, alpha=alpha, line_height=lh, letter_spacing=ls, align=align, anchor=anchor, caps=caps, name=name, para_space_after=para_after)


def _est_h(variants, size, width, lh, secondary_size=None) -> float:
    total = 0.0
    for text, secondary in variants:
        sz = secondary_size if (secondary and secondary_size) else size
        total += estimate_height(joined_plain([(text, False)]), sz, width, lh)
    return total


def _footer(ctx: Ctx) -> None:
    f = theme.FOOTER
    y = FOOTER_TOP
    logo_variant = "knockout" if ctx.dark else "small"
    base_color = theme.PAPER if ctx.dark else theme.INK
    base_alpha = 0.7 if ctx.dark else 0.55
    logo_w = 0.0
    if ctx.show_logo:
        pic = add_logo(ctx.slide, logo_variant, f["left"], y, f["logo_h"], name=f"logo:{logo_variant}")
        logo_w = pic.width / theme.EMU_PER_PX
    label = ctx.footer_label
    if label and ctx.lang == "bi":
        label = label[:1]  # PageFooter has no bilingual form; it shows the primary language only
    if label:
        lx = f["left"] + logo_w + (f["gap"] if logo_w else 0)
        _text(ctx, lx, y, 700, f["logo_h"], label, size=f["font"], color=base_color, alpha=base_alpha * 0.7, ls=0.4, anchor="middle", lh=1.0, name="footer:label")
    num = f"{ctx.page_number:02d}"
    add_text(ctx.slide, theme.ARTBOARD_W - f["right"] - 200, y, 200, f["logo_h"], [[{"text": num}]], size=f["font"], weight=600, color=base_color, alpha=base_alpha, letter_spacing=0.4, align="right", anchor="middle", line_height=1.0, name="footer:page")


def _eyebrow(ctx: Ctx, s: dict, *, y=70, w=CONTENT_W, weight=700, ls=2.5, color=None):
    v = ctx.copy(s, "eyebrow")
    if v:
        _text(ctx, MARGIN, y, w, 40, v, size=16, weight=weight, color=color or ctx.accent, ls=ls, caps=True, lh=1.2, name="text:eyebrow")


def _content_header(ctx: Ctx, s: dict, *, title_w=CONTENT_W, title_color=None):
    """Eyebrow (70) + title (110, 56px light) + accent rule (290) - the content grammar."""
    _eyebrow(ctx, s, w=title_w)
    title = ctx.copy(s, "title", required=True)
    color = title_color or theme.INK
    est = _est_h(title, 56, title_w, 1.05)
    if est > 180:
        ctx.warn(f"title likely exceeds 3 lines ({est:.0f}px) and would collide with the accent rule at y=290; shorten it")
    _text(ctx, MARGIN, 110, title_w, 180, title, size=56, weight=300, color=color, lh=1.05, ls=-1, name="text:title")
    add_rect(ctx.slide, MARGIN, 290, 60, 3, fill=ctx.accent, name="deco:rule")


def _meta_lines(ctx: Ctx, s: dict) -> list[tuple[str, bool]]:
    """`meta` is a list of short lines (location / date) or a single string."""
    meta = s.get("meta")
    if meta is None:
        return []
    if isinstance(meta, (str, dict)):
        return resolve_copy(meta, ctx.lang, path=f"slide {ctx.index + 1}.meta")
    out: list[tuple[str, bool]] = []
    for i, m in enumerate(meta):
        out.extend(resolve_copy(m, ctx.lang, path=f"slide {ctx.index + 1}.meta[{i}]"))
    return out


# ---------------------------------------------------------------------------
# COVERS
# ---------------------------------------------------------------------------


def cover1(ctx: Ctx, s: dict) -> None:
    """Hero photo on dark - the default cover."""
    _set_background(ctx.slide, theme.BLACK)
    photo, descr = _resolve_photo(ctx, s)
    cx, cy, cw, ch = theme.ARTBOARD_W - 60 - 720, 100, 720, 880
    add_gradient_rect(ctx.slide, cx, cy, cw, ch, stops=[(0, ctx.accent, None), (1, theme.DEEP_GRADIENT_END, None)], css_angle=135, radius=18, name="bg:photo-underlay")
    add_picture_cover(ctx.slide, photo, cx, cy, cw, ch, radius=18, alpha=0.92, name="photo:hero", descr=descr)
    add_gradient_rect(ctx.slide, cx - 240, 0, 240, theme.ARTBOARD_H, stops=[(0, theme.BLACK, 0.0), (1, theme.BLACK, 1.0)], css_angle=90, name="deco:vignette")

    title = ctx.copy(s, "title", required=True)
    tw = 900
    y = 220
    th = _est_h(title, 88, tw, 0.95)
    _text(ctx, MARGIN, y, tw, th, title, size=88, weight=300, color=theme.CREAM, lh=0.95, ls=-1.5, accent_color=theme.CYAN, name="text:title")
    y += th + 64
    sub = ctx.copy(s, "subtitle")
    if sub:
        sh = _est_h(sub, 28, tw, 1.25)
        _text(ctx, MARGIN, y, tw, sh, sub, size=28, color=theme.PAPER, alpha=0.86, lh=1.25, name="text:subtitle")
        y += sh + 90
    meta = _meta_lines(ctx, s)
    if meta:
        mh = len(meta) * 22
        _text(ctx, MARGIN, y, tw, mh, [(m, sec) for m, sec in meta], size=16, color=theme.PAPER, alpha=0.78, ls=0.4, lh=22 / 16, name="text:meta")
        y += mh
    logo_top = theme.ARTBOARD_H - 80 - 56
    if y > logo_top - theme.GROUP_GAP:
        raise SpecError(
            f"slide {ctx.index + 1} (cover1): the title / subtitle / meta block flows to y={y:.0f}px and would collide with the logo at y={logo_top}px "
            f"(keep {theme.GROUP_GAP}px clear). Shorten the title or subtitle, drop a meta line, or use a single language."
        )
    if ctx.show_logo:
        add_logo(ctx.slide, "knockout", MARGIN, theme.ARTBOARD_H - 80, 56, anchor="bl")


def cover2(ctx: Ctx, s: dict) -> None:
    """Editorial light cover - cream ground, big photo card right."""
    _set_background(ctx.slide, theme.CREAM)
    photo, descr = _resolve_photo(ctx, s)
    cx, cy, cw, ch = theme.ARTBOARD_W - 60 - 820, 60, 820, 960
    under = add_rect(ctx.slide, cx, cy, cw, ch, fill=ctx.accent, radius=24, name="bg:photo-underlay")
    add_shadow(under, blur_px=80, dist_px=30, css_deg=90, color=theme.INK, alpha=0.4)
    add_picture_cover(ctx.slide, photo, cx, cy, cw, ch, radius=24, name="photo:hero", descr=descr)

    add_rect(ctx.slide, MARGIN, 220, 80, 6, fill=ctx.accent, radius=3, name="deco:rule")
    _eyebrow(ctx, s, y=250, w=880, weight=600, ls=3)
    title = ctx.copy(s, "title", required=True)
    th = _est_h(title, 96, 880, 0.96)
    if th > 340:
        ctx.warn("cover title exceeds ~3 lines at 96px and will run into the subtitle at y=640")
    _text(ctx, MARGIN, 300, 880, 340, title, size=96, weight=300, color=theme.INK, lh=0.96, ls=-2, emphasis_weight=400, name="text:title")
    sub = ctx.copy(s, "subtitle")
    if sub:
        _text(ctx, MARGIN, 640, 880, 260, sub, size=26, color=theme.INK, alpha=0.78, lh=1.35, name="text:subtitle")
    bottom = theme.ARTBOARD_H - 80
    if ctx.show_logo:
        add_logo(ctx.slide, "primary", MARGIN, bottom, 56, anchor="bl")
    meta = _meta_lines(ctx, s)
    if meta:
        mh = len(meta) * 24
        _text(ctx, MARGIN + 380, bottom - mh, 500, mh, meta, size=16, color=theme.INK, alpha=0.7, ls=0.4, lh=1.5, align="right", anchor="bottom", name="text:meta")


def cover3(ctx: Ctx, s: dict) -> None:
    """Typographic cover - no image, diagonal hairline pattern, accent corner block."""
    _set_background(ctx.slide, theme.INK)
    add_diagonal_hairlines(ctx.slide, spacing=40, color=theme.PAPER, alpha=0.06, width_px=0.5)
    add_rect(ctx.slide, theme.ARTBOARD_W - 380, 0, 380, 380, fill=ctx.accent, name="deco:corner")
    number = s.get("number")
    if number is not None:
        add_text(ctx.slide, theme.ARTBOARD_W - 60 - 600, 320, 600, 220, [[{"text": str(number)}]], size=200, weight=200, color=ctx.accent, alpha=0.18, line_height=1.0, align="right", name="deco:ghost-number")
    _eyebrow(ctx, s, y=140, w=1350, weight=600, ls=3)
    title = ctx.copy(s, "title", required=True)
    th = _est_h(title, 120, 1350, 0.96)
    if th > 600:
        ctx.warn("cover title exceeds ~5 lines at 120px; shorten it")
    _text(ctx, MARGIN, 260, 1350, max(th, 240), title, size=120, weight=300, color=theme.CREAM, lh=0.96, ls=-2, accent_italic=True, accent_weight=400, name="text:title")
    bottom = theme.ARTBOARD_H - 80
    if ctx.show_logo:
        add_logo(ctx.slide, "knockout", MARGIN, bottom, 56, anchor="bl")
    meta = _meta_lines(ctx, s)
    if meta:
        mh = len(meta) * 24
        _text(ctx, theme.ARTBOARD_W - MARGIN - 700, bottom - mh, 700, mh, meta, size=16, color=theme.CREAM, alpha=0.7, ls=0.4, lh=1.5, align="right", anchor="bottom", name="text:meta")


# ---------------------------------------------------------------------------
# DIVIDERS
# ---------------------------------------------------------------------------


def _number(ctx: Ctx, s: dict) -> str:
    n = s.get("number")
    if n is None:
        raise SpecError(f"slide {ctx.index + 1} ({ctx.template}): `number` is required (e.g. '02')")
    return f"{int(n):02d}" if isinstance(n, int) else str(n)


def divider_image(ctx: Ctx, s: dict) -> None:
    _set_background(ctx.slide, theme.INK)
    photo, descr = _resolve_photo(ctx, s)
    add_picture_cover(ctx.slide, photo, theme.ARTBOARD_W - 60 - 1100, 100, 1100, 880, radius=24, name="photo:card", descr=descr)
    add_text(ctx.slide, 80, 360, 600, 220, [[{"text": _number(ctx, s)}]], size=220, weight=200, color=ctx.accent, line_height=1.0, letter_spacing=-8, name="text:number")
    title = ctx.copy(s, "title", required=True)
    _text(ctx, MARGIN, 600, 580, 100, title, size=64, weight=300, color=theme.CREAM, lh=1.0, ls=-1, name="text:title")
    cap = ctx.copy(s, "caption")
    if cap:
        ch = _est_h(cap, 22, 580, 1.4)
        if 700 + ch > BODY_LIMIT:
            ctx.warn("caption is too long for the divider; it would reach the footer band")
        _text(ctx, MARGIN, 700, 580, BODY_LIMIT - 700, cap, size=22, color=theme.CREAM, alpha=0.7, lh=1.4, name="text:caption")
    _footer(ctx)


def divider_dark(ctx: Ctx, s: dict) -> None:
    _set_background(ctx.slide, theme.INK)
    add_rect(ctx.slide, 0, 0, 12, theme.ARTBOARD_H, fill=ctx.accent, name="deco:bar")
    add_text(ctx.slide, MARGIN, 380, 460, 260, [[{"text": _number(ctx, s)}]], size=260, weight=200, color=ctx.accent, alpha=0.95, line_height=1.0, letter_spacing=-10, name="text:number")
    title = ctx.copy(s, "title", required=True)
    th = _est_h(title, 84, 1250, 1.0)
    if th > 168:
        ctx.warn("divider title exceeds 2 lines at 84px and would collide with the body copy at y=580")
    _text(ctx, 580, 410, 1250, 170, title, size=84, weight=300, color=theme.CREAM, lh=1.0, ls=-2, name="text:title")
    body = ctx.copy(s, "caption")
    if body:
        _text(ctx, 580, 580, 720, BODY_LIMIT - 580, body, size=26, color=theme.CREAM, alpha=0.7, lh=1.4, name="text:caption")
    _footer(ctx)


def divider_bright(ctx: Ctx, s: dict) -> None:
    _set_background(ctx.slide, theme.CREAM)
    add_rect(ctx.slide, theme.ARTBOARD_W - 720, theme.ARTBOARD_H - 720, 720, 720, fill=ctx.accent, round_corner="tl", radius=24, name="deco:field")
    _eyebrow(ctx, s, y=140, w=1000, weight=700, ls=3)
    add_text(ctx.slide, MARGIN, 220, 1000, 300, [[{"text": _number(ctx, s)}]], size=320, weight=200, color=theme.INK, line_height=0.9, letter_spacing=-12, name="text:number")
    title = ctx.copy(s, "title", required=True)
    _text(ctx, MARGIN, 620, 1000, 100, title, size=72, weight=300, color=theme.INK, lh=1.0, ls=-1.5, name="text:title")
    body = ctx.copy(s, "caption")
    if body:
        _text(ctx, MARGIN, 730, 800, BODY_LIMIT - 730, body, size=24, color=theme.INK, alpha=0.7, lh=1.4, name="text:caption")
    _footer(ctx)


# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------


def _body_paragraphs(ctx: Ctx, s: dict, key: str = "body") -> list[list[tuple[str, bool]]]:
    body = s.get(key)
    if body is None:
        return []
    items = body if isinstance(body, list) else [body]
    return [resolve_copy(b, ctx.lang, path=f"slide {ctx.index + 1}.{key}[{i}]") for i, b in enumerate(items)]


def _body_block(ctx: Ctx, x, y, w, h, paragraphs, *, size, color, alpha, lh, para_after, name):
    if not paragraphs:
        return
    flat: list[tuple[str, bool]] = []
    for p in paragraphs:
        flat.extend(p)
    est = sum(_est_h(p, size, w, lh) + para_after for p in paragraphs)
    if est > h:
        ctx.warn(f"{name}: ~{est:.0f}px of copy for a {h:.0f}px zone; trim the copy or split the slide")
    _text(ctx, x, y, w, h, flat, size=size, color=color, alpha=alpha, lh=lh, para_after=para_after, name=name)


def content_image_right(ctx: Ctx, s: dict) -> None:
    _set_background(ctx.slide, theme.PAPER)
    photo, descr = _resolve_photo(ctx, s)
    add_picture_cover(ctx.slide, photo, theme.ARTBOARD_W - 820, 0, 820, theme.ARTBOARD_H, name="photo:half", descr=descr)
    _content_header(ctx, s, title_w=880, title_color=ctx.accent)
    foot = ctx.copy(s, "footnote")
    body_bottom = 900 if foot else BODY_LIMIT
    _body_block(ctx, MARGIN, 340, 880, body_bottom - 340, _body_paragraphs(ctx, s), size=22, color=theme.BODY, alpha=0.85, lh=1.55, para_after=24, name="text:body")
    if foot:
        _text(ctx, MARGIN, 900, 880, 70, foot, size=14, color=theme.BODY, alpha=0.5, lh=1.4, anchor="bottom", name="text:footnote")
    _footer(ctx)


def content_columns(ctx: Ctx, s: dict) -> None:
    """Numbered columns (2-4 in one row, 5-6 in two rows) - the pillar / principle / step workhorse."""
    _set_background(ctx.slide, theme.PAPER)
    _content_header(ctx, s)
    cols = s.get("columns")
    if not cols or not isinstance(cols, list):
        raise SpecError(f"slide {ctx.index + 1} (content_columns): `columns` (list of {{number, heading, body}}) is required")
    n = len(cols)
    if n < 2 or n > 6:
        raise SpecError(f"slide {ctx.index + 1} (content_columns): 2 to 6 columns are supported, got {n}")
    per_row = n if n <= 4 else 3
    rows = -(-n // per_row)
    gap = 48
    col_w = (CONTENT_W - (per_row - 1) * gap) / per_row
    row_h = (BODY_LIMIT - 340 - (rows - 1) * gap) / rows
    for i, c in enumerate(cols):
        r, k = divmod(i, per_row)
        x = MARGIN + k * (col_w + gap)
        top = 340 + r * (row_h + gap)
        number = c.get("number", f"{i + 1:02d}")
        number = f"{int(number):02d}" if isinstance(number, int) else str(number)
        add_text(ctx.slide, x, top, col_w, 64, [[{"text": number}]], size=64, weight=200, color=ctx.accent, line_height=1.0, letter_spacing=-2, name=f"text:col{i + 1}-number")
        heading = resolve_copy(c.get("heading"), ctx.lang, path=f"slide {ctx.index + 1}.columns[{i}].heading")
        if not heading:
            raise SpecError(f"slide {ctx.index + 1} (content_columns): columns[{i}].heading is required")
        hy = top + 64 + 24
        hh = _est_h(heading, 26, col_w, 1.2, secondary_size=22)
        _text(ctx, x, hy, col_w, hh, heading, size=26, weight=600, color=theme.INK, lh=1.2, secondary_size=22, name=f"text:col{i + 1}-heading")
        body = resolve_copy(c.get("body"), ctx.lang, path=f"slide {ctx.index + 1}.columns[{i}].body")
        if body:
            by = hy + hh + 16
            bh = top + row_h - by
            est = _est_h(body, 18, col_w, 1.55)
            if est > bh:
                ctx.warn(f"columns[{i}] body needs ~{est:.0f}px but only {bh:.0f}px remain; trim it")
            _text(ctx, x, by, col_w, bh, body, size=18, color=theme.BODY, alpha=0.78, lh=1.55, name=f"text:col{i + 1}-body")
    _footer(ctx)


def content_stat(ctx: Ctx, s: dict) -> None:
    """One hero number with a supporting metric table."""
    _set_background(ctx.slide, theme.CREAM)
    _content_header(ctx, s)
    stat = s.get("stat")
    if not stat or "value" not in stat:
        raise SpecError(f"slide {ctx.index + 1} (content_stat): `stat: {{value, unit?, caption}}` is required")
    left_w = (CONTENT_W - 80) * 1.3 / 2.3
    right_x = MARGIN + left_w + 80
    right_w = CONTENT_W - 80 - left_w
    spans = [{"text": str(stat["value"]), "color": ctx.accent}]
    if stat.get("unit"):
        spans.append({"text": " " + str(stat["unit"]), "size": 120, "baseline": 75000})
    add_text(ctx.slide, MARGIN, 360, left_w, 234, [spans], size=260, weight=200, color=theme.INK, line_height=0.9, letter_spacing=-10, name="text:stat")
    cap = resolve_copy(stat.get("caption"), ctx.lang, path=f"slide {ctx.index + 1}.stat.caption")
    if cap:
        _text(ctx, MARGIN, 618, min(700, left_w), 120, cap, size=26, color=theme.INK, lh=1.3, name="text:stat-caption")
    rows = s.get("rows") or []
    if len(rows) > 6:
        raise SpecError(f"slide {ctx.index + 1} (content_stat): at most 6 supporting rows fit; got {len(rows)}")
    y = 400
    row_h = 82
    for i, r in enumerate(rows):
        label = resolve_copy(r.get("label"), ctx.lang, path=f"slide {ctx.index + 1}.rows[{i}].label")
        value = str(r.get("value", ""))
        if not label or not value:
            raise SpecError(f"slide {ctx.index + 1} (content_stat): rows[{i}] needs `label` and `value`")
        baseline_y = y + 22 + 38
        _text(ctx, right_x, baseline_y - 38, right_w - 260, 38, label, size=20, color=theme.INK, alpha=0.7, lh=1.0, anchor="bottom", name=f"text:row{i + 1}-label")
        add_text(ctx.slide, right_x + right_w - 250, baseline_y - 38, 250, 38, [[{"text": value}]], size=32, weight=500, color=theme.INK, line_height=1.0, align="right", anchor="bottom", name=f"text:row{i + 1}-value")
        add_hairline(ctx.slide, right_x, y + row_h - 1, right_w, color=theme.INK, alpha=0.12)
        y += row_h
    _footer(ctx)


def content_chart(ctx: Ctx, s: dict) -> None:
    """Content grammar + a native chart (teal tints only, never a pie)."""
    _set_background(ctx.slide, theme.PAPER)
    _content_header(ctx, s)
    chart = s.get("chart")
    if not chart:
        raise SpecError(f"slide {ctx.index + 1} (content_chart): `chart` is required")
    paragraphs = _body_paragraphs(ctx, s)
    top, h = 340, BODY_LIMIT - 340
    if paragraphs:
        _body_block(ctx, MARGIN, top, 560, h, paragraphs, size=22, color=theme.BODY, alpha=0.85, lh=1.55, para_after=24, name="text:body")
        add_chart(ctx.slide, MARGIN + 560 + 48, top, CONTENT_W - 560 - 48, h, chart)
    else:
        add_chart(ctx.slide, MARGIN, top, CONTENT_W, h, chart)
    if s.get("footnote"):
        ctx.warn("`footnote` is not part of the chart template; put the source in `footer_label` instead")
    _footer(ctx)


def content_table(ctx: Ctx, s: dict) -> None:
    _set_background(ctx.slide, theme.PAPER)
    _content_header(ctx, s)
    table = s.get("table")
    if not table:
        raise SpecError(f"slide {ctx.index + 1} (content_table): `table: {{headers, rows}}` is required")
    headers = [joined_plain(resolve_copy(hd, ctx.lang, path=f"slide {ctx.index + 1}.table.headers[{i}]")) for i, hd in enumerate(table.get("headers", []))]
    rows = [[joined_plain(resolve_copy(c, ctx.lang, path=f"slide {ctx.index + 1}.table.rows[{ri}][{ci}]")) if isinstance(c, dict) else c for ci, c in enumerate(r)] for ri, r in enumerate(table.get("rows", []))]
    _, h = add_table(ctx.slide, MARGIN, 340, CONTENT_W, headers, rows)
    if 340 + h > BODY_LIMIT:
        raise SpecError(f"slide {ctx.index + 1} (content_table): {len(rows)} rows do not fit above the footer band; at most {int((BODY_LIMIT - 340 - 56) // 64)} rows fit")
    _footer(ctx)


def back_cover(ctx: Ctx, s: dict) -> None:
    """Plain closing slide: cream surface, centred lockup, no copy, no page number."""
    _set_background(ctx.slide, theme.CREAM)
    add_logo(ctx.slide, "primary", theme.ARTBOARD_W / 2, theme.ARTBOARD_H / 2, 96, anchor="c")


TEMPLATES = {
    "cover1": cover1,
    "cover2": cover2,
    "cover3": cover3,
    "divider_image": divider_image,
    "divider_dark": divider_dark,
    "divider_bright": divider_bright,
    "content_image_right": content_image_right,
    "content_columns": content_columns,
    "content_stat": content_stat,
    "content_chart": content_chart,
    "content_table": content_table,
    "back_cover": back_cover,
}
COVER_TEMPLATES = ("cover1", "cover2", "cover3")
