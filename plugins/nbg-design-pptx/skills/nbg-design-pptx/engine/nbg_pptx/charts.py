"""Native chart and table components styled with the NBG design-system tokens.

The HTML design system has no chart template; these components live inside the
content-slide grammar (eyebrow / title / rule) and follow its rules: teal tints only,
no pie charts (doughnut instead), calm monochrome frames, Aptos throughout.
"""

from __future__ import annotations

import re

from lxml import etree
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from . import theme
from .primitives import _apply_run_style, _srgb
from .text import SpecError
from .theme import px

_CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
}

AXIS_INK = "5A5F5A"  # ~ rgba(0,56,65,0.7) on a light surface


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def add_chart(slide, x: float, y: float, w: float, h: float, spec: dict, *, name: str = "chart:main", palette: tuple[str, ...] | None = None, axis_ink: str = AXIS_INK, grid: str = theme.GREY_1, label_ink: str = theme.INK):
    ctype = spec.get("type")
    if ctype is None:
        raise SpecError("chart.type is required (column | bar | line | area | doughnut | stacked_column | stacked_bar)")
    if ctype == "pie":
        raise SpecError("chart.type 'pie' is not allowed by the NBG design guidelines; use 'doughnut'")
    if ctype not in _CHART_TYPES:
        raise SpecError(f"chart.type {ctype!r} is not supported; use one of {sorted(_CHART_TYPES)}")
    categories = spec.get("categories")
    series = spec.get("series")
    if not categories or not series:
        raise SpecError("chart.categories and chart.series are required")

    data = CategoryChartData()
    data.categories = list(categories)
    for i, s in enumerate(series):
        if "name" not in s or "values" not in s:
            raise SpecError(f"chart.series[{i}] needs `name` and `values`")
        if len(s["values"]) != len(categories):
            raise SpecError(f"chart.series[{i}] has {len(s['values'])} values for {len(categories)} categories")
        data.add_series(s["name"], list(s["values"]))

    gf = slide.shapes.add_chart(_CHART_TYPES[ctype], px(x), px(y), px(w), px(h), data)
    gf.name = name
    chart = gf.chart
    chart.has_title = False
    chart.font.name = theme.FONT_FAMILY
    chart.font.size = Pt(theme.pt(18))
    chart.font.color.rgb = _rgb(axis_ink)

    multi = len(series) > 1
    chart.has_legend = multi
    if multi:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(theme.pt(18))
        chart.legend.font.color.rgb = _rgb(axis_ink)

    plot = chart.plots[0]
    number_format = spec.get("number_format", "General")
    show_values = spec.get("data_labels", True)
    if show_values:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(theme.pt(18))
        dl.font.color.rgb = _rgb(label_ink)
        dl.number_format = number_format
        dl.number_format_is_linked = False
        if ctype in ("column", "bar"):
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        elif ctype in ("stacked_column", "stacked_bar"):
            dl.position = XL_LABEL_POSITION.CENTER
        elif ctype == "doughnut":
            dl.show_percentage = spec.get("show_percentage", True)
            dl.show_value = not spec.get("show_percentage", True)

    palette = list(palette or theme.CHART_SERIES)
    if ctype == "doughnut":
        # colour points, not series
        s0 = plot.series[0]
        for i, _ in enumerate(categories):
            pt_ = s0.points[i]
            pt_.format.fill.solid()
            pt_.format.fill.fore_color.rgb = _rgb(palette[i % len(palette)])
            pt_.format.line.fill.background()
        hole = plot._element.find(qn("c:holeSize"))
        if hole is None:
            hole = etree.SubElement(plot._element, qn("c:holeSize"))
        hole.set("val", str(int(spec.get("hole_size", 60))))
    else:
        for i, s in enumerate(plot.series):
            colour = palette[i % len(palette)]
            if ctype in ("line",):
                s.format.line.color.rgb = _rgb(colour)
                s.format.line.width = Pt(3.5)
                s.smooth = False
                s.marker.style = XL_MARKER_STYLE.CIRCLE
                s.marker.size = 9
                s.marker.format.fill.solid()
                s.marker.format.fill.fore_color.rgb = _rgb(theme.PAPER)
                s.marker.format.line.color.rgb = _rgb(colour)
                s.marker.format.line.width = Pt(3.5)
            else:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = _rgb(colour)
                s.format.line.fill.background()
                if ctype == "area":
                    # subtle fill, like the HTML area-line treatment
                    srgb = s.format.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
                    a = etree.SubElement(srgb, qn("a:alpha"))
                    a.set("val", "18000")
        if ctype in ("column", "bar", "stacked_column", "stacked_bar"):
            plot.gap_width = int(spec.get("gap_width", 80))
            if ctype in ("stacked_column", "stacked_bar"):
                plot.overlap = 100

        # Calm frame: light value gridlines, no axis lines, muted labels.
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = _rgb(grid)
        va.major_gridlines.format.line.width = Pt(0.5)
        va.format.line.fill.background()
        va.tick_labels.font.size = Pt(theme.pt(16))
        va.tick_labels.font.color.rgb = _rgb(axis_ink)
        va.has_title = False
        if spec.get("hide_value_axis", False):
            va.visible = False
        ca = chart.category_axis
        ca.format.line.fill.background()
        ca.tick_labels.font.size = Pt(theme.pt(18))
        ca.tick_labels.font.color.rgb = _rgb(axis_ink)
        ca.has_major_gridlines = False
    return gf


# ---------------------------------------------------------------------------
# Table
# ---------------------------------------------------------------------------

_NUMERIC = re.compile(r"^\s*[+\-−]?[\d.,\s]+(%|[kKmMbB]|bn|pp|bps|€|\$|£)?\s*$|^\s*[€$£]\s*[\d.,\s]+[kKmMbB]?\s*$")


def _ln(tag: str, *, color: str | None, alpha: float | None = None, width_px: float = 0):
    ln = etree.Element(qn(tag))
    if color is None:
        ln.set("w", "0")
        etree.SubElement(ln, qn("a:noFill"))
    else:
        ln.set("w", str(px(width_px)))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sf.append(_srgb(color, alpha))
    return ln


def _set_cell_borders(cell, *, bottom_color: str | None, bottom_alpha: float | None, bottom_px: float):
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tc_pr.findall(qn(tag)):
            tc_pr.remove(el)
    borders = [
        _ln("a:lnL", color=None),
        _ln("a:lnR", color=None),
        _ln("a:lnT", color=None),
        _ln("a:lnB", color=bottom_color, alpha=bottom_alpha, width_px=bottom_px),
    ]
    for i, b in enumerate(borders):
        tc_pr.insert(i, b)


def add_table(slide, x: float, y: float, w: float, headers: list[str], rows: list[list], *, name: str = "table:main"):
    if not headers or not rows:
        raise SpecError("table.headers and table.rows are required")
    n_cols = len(headers)
    for i, r in enumerate(rows):
        if len(r) != n_cols:
            raise SpecError(f"table.rows[{i}] has {len(r)} cells for {n_cols} headers")
    header_h, row_h = 56, 64
    h = header_h + row_h * len(rows)
    gf = slide.shapes.add_table(len(rows) + 1, n_cols, px(x), px(y), px(w), px(h))
    gf.name = name
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tbl.vert_banding = False
    tbl.last_row = False
    tbl.first_col = False
    tbl.last_col = False

    first_w = 0.34 * w if n_cols > 1 else w
    other_w = (w - first_w) / max(1, n_cols - 1)
    for ci in range(n_cols):
        tbl.columns[ci].width = px(first_w if ci == 0 else other_w)
    tbl.rows[0].height = px(header_h)
    for ri in range(1, len(rows) + 1):
        tbl.rows[ri].height = px(row_h)

    def style_cell(cell, text: str, *, header: bool, numeric: bool, last: bool):
        cell.fill.background()
        cell.margin_left = 0
        cell.margin_right = px(16)
        cell.margin_top = px(12)
        cell.margin_bottom = px(12)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if numeric and not header else (PP_ALIGN.RIGHT if numeric else PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text)
        if header:
            _apply_run_style(run, size_px=18, weight=600, color=theme.INK, alpha=None, italic=False, letter_spacing=0, caps=False)
        elif numeric:
            _apply_run_style(run, size_px=22, weight=500, color=theme.INK, alpha=None, italic=False, letter_spacing=0, caps=False)
        else:
            _apply_run_style(run, size_px=20, weight=400, color=theme.BODY, alpha=0.85, italic=False, letter_spacing=0, caps=False)
        if header:
            _set_cell_borders(cell, bottom_color=theme.INK, bottom_alpha=None, bottom_px=2)
        else:
            _set_cell_borders(cell, bottom_color=theme.INK, bottom_alpha=0.12, bottom_px=1)

    numeric_cols = [
        all(_NUMERIC.match(str(r[ci])) for r in rows if str(r[ci]).strip() not in ("", "-", "—"))
        for ci in range(n_cols)
    ]
    for ci, htext in enumerate(headers):
        style_cell(tbl.cell(0, ci), htext, header=True, numeric=numeric_cols[ci] and ci > 0, last=False)
    for ri, r in enumerate(rows, start=1):
        for ci, val in enumerate(r):
            style_cell(tbl.cell(ri, ci), val, header=False, numeric=numeric_cols[ci] and ci > 0, last=ri == len(rows))
    return gf, h
