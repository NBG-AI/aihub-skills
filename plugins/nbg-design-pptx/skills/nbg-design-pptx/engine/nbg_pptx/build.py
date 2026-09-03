"""Deck spec (YAML) -> native NBG-design PPTX."""

from __future__ import annotations

from pathlib import Path

import yaml
from pptx import Presentation

from . import theme
from .briefing import BRIEF_COVERS, BRIEF_TEMPLATES
from .templates import COVER_TEMPLATES, TEMPLATES as DESIGN_TEMPLATES, Ctx
from .text import SpecError, resolve_copy

# The only approved defaults (SKILL.md "Approved defaults"): language en, logo shown,
# style briefing. Everything else missing is an error, never a silent fallback.
APPROVED_DEFAULTS = {"language": "en", "show_logo": True, "style": theme.DEFAULT_STYLE}

FAMILIES = {
    "briefing": {"templates": BRIEF_TEMPLATES, "covers": BRIEF_COVERS},
    "design-system": {"templates": DESIGN_TEMPLATES, "covers": COVER_TEMPLATES},
}
TEMPLATES = {**DESIGN_TEMPLATES, **BRIEF_TEMPLATES}


def load_spec(path: Path) -> dict:
    with open(path, encoding="utf-8") as f:
        spec = yaml.safe_load(f)
    if not isinstance(spec, dict) or "slides" not in spec:
        raise SpecError("spec must be a mapping with a `slides` list")
    return spec


def _accent_for(template: str, deck_accent: str | None, slide_accent: str | None, index: int) -> str:
    raw = slide_accent or deck_accent
    if raw is None:
        return theme.TEMPLATE_DEFAULT_ACCENT[template]
    hex6 = theme.normalize_hex(str(raw))
    if hex6 not in theme.ACCENTS:
        raise SpecError(f"slide {index + 1}: accent #{hex6} is outside the four allowed accents {['#' + a for a in theme.ACCENTS]}")
    surface = theme.TEMPLATE_SURFACE[template]
    if theme.contrast_ratio(hex6, surface) < 1.6:
        raise SpecError(f"slide {index + 1} ({template}): accent #{hex6} is invisible on the template surface #{surface}; choose another accent for this slide")
    return hex6


def build(spec_path: Path, out_path: Path) -> list[str]:
    spec = load_spec(spec_path)
    deck = spec.get("deck") or {}
    lang = deck.get("language", APPROVED_DEFAULTS["language"])
    if lang not in ("en", "gr", "bi"):
        raise SpecError(f"deck.language must be en | gr | bi, got {lang!r}")
    show_logo = bool(deck.get("show_logo", APPROVED_DEFAULTS["show_logo"]))
    style = deck.get("style", APPROVED_DEFAULTS["style"])
    if style not in theme.STYLES:
        raise SpecError(f"deck.style must be one of {list(theme.STYLES)}, got {style!r}")
    family = FAMILIES[style]
    deck_accent = deck.get("accent")
    if style == "briefing" and deck_accent is not None:
        raise SpecError("deck.accent is not configurable in the briefing style (the accent is fixed to #007B85)")
    slides = spec["slides"]
    if not isinstance(slides, list) or not slides:
        raise SpecError("`slides` must be a non-empty list")

    prs = Presentation()
    prs.slide_width = theme.SLIDE_W_EMU
    prs.slide_height = theme.SLIDE_H_EMU
    blank = prs.slide_layouts[6]
    if deck.get("title"):
        prs.core_properties.title = str(deck["title"])
    prs.core_properties.author = "National Bank of Greece"
    lean = deck.get("lean", True)
    if lean not in (True, False):
        raise SpecError("deck.lean must be true or false (false = a full-length deck; the lean limits become warnings)")
    prs.core_properties.category = f"nbg-design-pptx:{style}" + ("" if lean else ";lean=off")

    warnings: list[str] = []
    default_footer = deck.get("footer_label")
    for i, s in enumerate(slides):
        if not isinstance(s, dict) or "template" not in s:
            raise SpecError(f"slide {i + 1}: each slide needs a `template` (one of {sorted(TEMPLATES)})")
        template = s["template"]
        if template not in TEMPLATES:
            raise SpecError(f"slide {i + 1}: unknown template {template!r}; use one of {sorted(TEMPLATES)}")
        if template not in family["templates"]:
            raise SpecError(f"slide {i + 1}: template {template!r} belongs to the other style; deck.style is {style!r} whose templates are {sorted(family['templates'])}")
        if i == 0 and template not in family["covers"]:
            raise SpecError(f"the first slide must be a cover ({' | '.join(sorted(family['covers']))})")
        slide = prs.slides.add_slide(blank)
        slide.name = template
        footer_label = s.get("footer_label", default_footer)
        if style == "briefing":
            if s.get("accent") is not None:
                raise SpecError(f"slide {i + 1}: `accent` is not configurable in the briefing style")
            accent, surface = theme.BRIEF["accent"], theme.BRIEF["white"]
        else:
            accent, surface = _accent_for(template, deck_accent, s.get("accent"), i), theme.TEMPLATE_SURFACE[template]
        ctx = Ctx(
            slide=slide,
            lang=lang,
            template=template,
            accent=accent,
            surface=surface,
            show_logo=show_logo,
            page_number=i + 1,
            footer_label=resolve_copy(footer_label, lang, path=f"slide {i + 1}.footer_label") if footer_label else None,
            index=i,
            total=len(slides),
        )
        TEMPLATES[template](ctx, s)
        if s.get("notes"):
            notes = resolve_copy(s["notes"], lang, path=f"slide {i + 1}.notes")
            slide.notes_slide.notes_text_frame.text = "\n\n".join(t for t, _ in notes)
        warnings.extend(ctx.warnings)

    if style == "design-system" and slides[-1].get("template") != "back_cover":
        warnings.append("the deck does not end with a `back_cover` slide (plain closing slide with the centred lockup)")
    if style == "briefing" and len(slides) > theme.BRIEF_LEAN_WARN:
        warnings.append(f"{len(slides)} slides: a briefing deck is lean - aim for {theme.BRIEF_LEAN_WARN} or fewer" + ("" if not lean else f" (hard limit {theme.BRIEF_LEAN_FAIL})"))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return warnings
